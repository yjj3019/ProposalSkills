"""e2e 골든 — 미니 RFP → slides.json → build_deck → deck_check → quality_gate → meta → audit → unified_gate.

fixtures/e2e-mini-rfp 가 스킬 파이프라인 전 구간을 재현한다. LibreOffice가 없으면 렌더 단계만 건너뛴다.
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

REPO = Path(__file__).resolve().parent
DOC = REPO / "skills/create-proposal-document"
BEST = REPO / "skills/create-best-proposal"
FIX = DOC / "fixtures/e2e-mini-rfp"
PY = sys.executable

sys.path.insert(0, str(DOC / "scripts"))
import build_deck  # noqa: E402
import deck_check  # noqa: E402


def run(*args, **kw) -> subprocess.CompletedProcess:
    return subprocess.run([PY, *map(str, args)], capture_output=True, text=True, encoding="utf-8",
                          errors="replace", cwd=str(REPO), **kw)


class DeckPipelineTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.tmp = tempfile.TemporaryDirectory()
        cls.out = Path(cls.tmp.name)
        cls.pptx = cls.out / "draft.pptx"
        proc = run(DOC / "scripts/build_deck.py", FIX / "slides.json", "-o", cls.pptx, "--strict")
        assert proc.returncode == 0, proc.stdout + proc.stderr
        cls.build_out = proc.stdout

    @classmethod
    def tearDownClass(cls):
        cls.tmp.cleanup()

    # ---- build_deck -----------------------------------------------------------
    def test_build_covers_every_layout_type(self):
        spec = json.loads((FIX / "slides.json").read_text(encoding="utf-8"))
        used = {s["type"] for s in spec["slides"]}
        self.assertEqual(used | {"bullets"}, build_deck.ALL_TYPES, f"미사용 레이아웃: {build_deck.ALL_TYPES - used}")
        self.assertIn("17 slides", self.build_out)

    def test_roles_named_on_body_slides(self):
        from pptx import Presentation
        prs = Presentation(str(self.pptx))
        body = [s for s in prs.slides if deck_check.slide_kind(s) == "body"]
        self.assertGreaterEqual(len(body), 8)
        for s in body:
            names = {sh.name for sh in deck_check.iter_shapes(s.shapes)}
            self.assertTrue({"HEADER", "TITLE", "LEAD", "CAPTION", "REQID", "FOOTER", "PAGENO"} <= names, names)

    def test_matrix_auto_paginates_with_header_repeat(self):
        spec = {"meta": {"title": "t", "bidder": "b"},
                "slides": [{"type": "matrix", "title": "조견표", "lead": "L", "rows_per_slide": 5,
                            "rows": [{"id": f"R{i}", "text": "x", "support": "O", "response_loc": "p", "note": ""}
                                     for i in range(1, 13)]}]}
        p = self.out / "matrix.pptx"
        (self.out / "m.json").write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        proc = run(DOC / "scripts/build_deck.py", self.out / "m.json", "-o", p)
        self.assertEqual(proc.returncode, 0, proc.stderr)
        from pptx import Presentation
        prs = Presentation(str(p))
        self.assertEqual(len(prs.slides), 3)
        titles = [deck_check.shape_text(next(sh for sh in s.shapes if sh.name == "TITLE")) for s in prs.slides]
        self.assertIn("(1/3)", titles[0])
        self.assertIn("(계속 3/3)", titles[2])

    def test_strict_fails_on_missing_lead_and_page_limit(self):
        spec = {"meta": {"title": "t", "page_limit": 1},
                "slides": [{"type": "cover"}, {"type": "bullets", "title": "x", "items": ["a"]}]}
        (self.out / "bad.json").write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        proc = run(DOC / "scripts/build_deck.py", self.out / "bad.json", "-o", self.out / "bad.pptx", "--strict")
        self.assertEqual(proc.returncode, 1)
        self.assertIn("리드문 없음", proc.stdout)
        self.assertIn("페이지 제한", proc.stdout)

    def test_unknown_type_is_usage_error(self):
        (self.out / "u.json").write_text(json.dumps({"slides": [{"type": "hero"}]}), encoding="utf-8")
        proc = run(DOC / "scripts/build_deck.py", self.out / "u.json", "-o", self.out / "u.pptx")
        self.assertEqual(proc.returncode, 2)

    # ---- deck_check -----------------------------------------------------------
    def test_deck_check_passes_golden_and_counts_body_pages(self):
        proc = run(DOC / "scripts/deck_check.py", self.pptx, "--max-pages", "40", "--exclude-cover-toc",
                   "--require-req-ids", "--stage", "draft")
        self.assertEqual(proc.returncode, 0, proc.stdout)
        self.assertIn("본문 14장 / 제한 40장", proc.stdout)

    def test_deck_check_blocks_small_font_and_missing_lead(self):
        spec = {"meta": {"title": "t"},
                "slides": [{"type": "bullets", "title": "x", "lead": "", "items": ["a"]}]}
        (self.out / "nl.json").write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        run(DOC / "scripts/build_deck.py", self.out / "nl.json", "-o", self.out / "nl.pptx")
        proc = run(DOC / "scripts/deck_check.py", self.out / "nl.pptx")
        self.assertEqual(proc.returncode, 1)
        self.assertIn("리드문 없음", proc.stdout)
        proc = run(DOC / "scripts/deck_check.py", self.pptx, "--min-font", "12")
        self.assertEqual(proc.returncode, 1)
        self.assertIn("최소 폰트", proc.stdout)

    def test_deck_check_page_limit_blocks(self):
        proc = run(DOC / "scripts/deck_check.py", self.pptx, "--max-pages", "10")
        self.assertEqual(proc.returncode, 1)
        self.assertIn("> 페이지 제한 10장", proc.stdout)

    def test_deck_check_heuristic_on_external_deck(self):
        from pptx import Presentation
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_textbox(0, 0, 100000, 100000).text_frame.text = "제목만 있는 외부 덱"
        p = self.out / "ext.pptx"
        prs.save(str(p))
        proc = run(DOC / "scripts/deck_check.py", p)
        self.assertEqual(proc.returncode, 1)
        self.assertIn("[휴리스틱]", proc.stdout)

    @unittest.skipUnless(shutil.which("soffice") or shutil.which("libreoffice"), "LibreOffice 없음")
    def test_render_emits_verified_block(self):
        rj = self.out / "render.json"
        proc = run(DOC / "scripts/deck_check.py", self.pptx, "--render", "--emit-render", rj, "--max-pages", "40")
        self.assertEqual(proc.returncode, 0, proc.stdout)
        block = json.loads(rj.read_text(encoding="utf-8"))
        self.assertTrue(block["verified"])
        self.assertTrue(block["artifact_hash"].startswith("sha256:"))
        self.assertTrue(any("렌더 OK" in e for e in block["evidence"]))

    def test_emit_render_without_soffice_is_not_verified(self):
        rj = self.out / "render_none.json"
        env = {"PATH": str(self.out)}  # soffice 못 찾게 PATH 비움
        import os
        env.update({k: v for k, v in os.environ.items() if k not in {"PATH"}})
        proc = run(DOC / "scripts/deck_check.py", self.pptx, "--emit-render", rj, env=env)
        self.assertEqual(proc.returncode, 0, proc.stdout)
        block = json.loads(rj.read_text(encoding="utf-8"))
        self.assertFalse(block["verified"])
        self.assertTrue(any("NOT INSPECTED" in e for e in block["evidence"]))

    # ---- quality_gate on the built deck ----------------------------------------
    def test_quality_gate_draft_warns_submission_blocks(self):
        qg = DOC / "scripts/quality_gate.py"
        draft = run(qg, self.pptx, "--stage", "draft")
        self.assertEqual(draft.returncode, 0, draft.stdout)  # [NEEDS INPUT]·○○는 초안 경고
        self.assertIn("[검토필요]", draft.stdout)
        sub = run(qg, self.pptx, "--stage", "submission")
        self.assertEqual(sub.returncode, 1)
        self.assertIn("needs input", sub.stdout)

    # ---- meta → audit → unified_gate -------------------------------------------
    def test_mid_draft_meta_is_blocked_with_worklist(self):
        audit = self.out / "audit.json"
        proc = run(BEST / "scripts/build_audit_from_meta.py", FIX / "meta.json", "-o", audit)
        self.assertEqual(proc.returncode, 0, proc.stderr)
        data = json.loads(audit.read_text(encoding="utf-8"))
        self.assertFalse(data["artifact_required"])  # draft 기본값
        gate = run(BEST / "scripts/unified_gate.py", audit, "--no-explain")
        self.assertEqual(gate.returncode, 1)
        for needle in ("requirement R6 is not approved", "blocking input I1 is open",
                       "missing attachment: 유사실적증명서(2건)", "unresolved token"):
            self.assertIn(needle, gate.stdout)
        self.assertIn("STATUS: BLOCKED", gate.stdout)

    def test_closed_draft_reaches_draft_ready(self):
        meta = json.loads((FIX / "meta.json").read_text(encoding="utf-8"))
        meta["bid_decision"] = "bid"
        meta["bid_conditions"] = []
        meta["eligibility"][0]["met"] = True
        for r in meta["requirements"]:
            if r["id"] == "R6":
                r.update(state="approved", evidence_refs=["첨부:유사실적증명서 2건"], reviewer="영업")
        meta["attachments"][1]["present"] = True
        meta["inputs"][0]["status"] = "closed"
        meta["unresolved_tokens"] = []
        mp = self.out / "meta_ready.json"
        mp.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")
        audit = self.out / "audit_ready.json"
        self.assertEqual(run(BEST / "scripts/build_audit_from_meta.py", mp, "-o", audit, "--strict").returncode, 0)
        gate = run(BEST / "scripts/unified_gate.py", audit, "--no-explain")
        self.assertEqual(gate.returncode, 0, gate.stdout)
        self.assertIn("STATUS: DRAFT-READY", gate.stdout)
        self.assertNotIn("SUBMISSION-READY", gate.stdout)
        # 같은 audit을 submission으로 올리면 렌더·패키지·클리어 요구로 막혀야 한다
        meta["mode"] = "submission"
        mp.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")
        run(BEST / "scripts/build_audit_from_meta.py", mp, "-o", audit)
        gate = run(BEST / "scripts/unified_gate.py", audit, "--no-explain")
        self.assertEqual(gate.returncode, 1)
        self.assertIn("render verification is missing", gate.stdout)


if __name__ == "__main__":
    unittest.main()
