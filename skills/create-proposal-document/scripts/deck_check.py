#!/usr/bin/env python3
"""제안서 PPTX 레이아웃 린트 + 렌더 검사 (review-checklist 7·8번 항목의 기계화).

사용법:
    python3 deck_check.py 제안서.pptx [--max-pages 40] [--exclude-cover-toc]
        [--render] [--png-dir out/] [--emit-render render.json] [--min-font PT]
        [--require-req-ids] [--stage draft|submission] [--profile 규격]

레이아웃 린트(python-pptx, 항상 수행):
  1. 페이지 수 vs --max-pages           5. 본문 최소 폰트 미만 — 캡션·헤더·푸터 제외
  2. 본문 장표의 리드문(LEAD) 존재·길이  6. 텍스트 과밀(표 셀 포함)
  3. 본문 장표의 REQ-ID 표기            7. 표 헤더 행 존재, 행 수 초과 표
  4. 제목(TITLE) 존재                    8. 빈 슬라이드·이미지 전용 슬라이드
  폰트·밀도·행 수 기준은 산출물 규격(deck_profiles.py)에서 온다. PPTX에 남은 규격 표시를
  읽어 자동 적용하며, --profile로 덮어쓸 수 있다. 표시가 없거나 모르는 값이면 제출
  단계에서 차단한다(가장 느슨한 기본값으로 조용히 통과시키지 않는다).
  build_deck.py 산출물은 도형 이름(TITLE/LEAD/REQID)으로 정확히 검사하고,
  외부 제작 덱은 위치 휴리스틱(상단 두 번째 텍스트 = 리드문)으로 검사한 뒤 [휴리스틱] 표기.

렌더 검사(--render, LibreOffice 필요):
  soffice → PDF 변환 → 페이지 수 대조 → (--png-dir) 장표별 PNG 썸네일.
  soffice가 없으면 차단하지 않고 "NOT INSPECTED"로 기록한다(통과 추정 금지).

--emit-render: audit용 render 블록을 JSON으로 기록한다
  {"verified": bool, "artifact_hash": "sha256:…", "tool": "…", "evidence": [...]}
  verified는 렌더 성공 + 차단 항목 0일 때만 true다.

종료 코드: 0=통과(경고 가능), 1=차단, 2=사용 오류·파일 손상.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import deck_profiles
import quality_gate

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:  # pragma: no cover
    print("python-pptx가 필요합니다: pip install python-pptx", file=sys.stderr)
    raise SystemExit(2)

# 기본값은 deck_profiles의 기본 프로파일에서 온다. 프로파일별 실제 기준은 lint()에
# 전달되며, 생성기와 같은 정의를 읽으므로 두 도구의 숫자가 갈리지 않는다.
_DEFAULT = deck_profiles.get(deck_profiles.DEFAULT_PROFILE)
LEAD_MAX = _DEFAULT["lead_max_chars"]
DENSITY_MAX = _DEFAULT["density_max"]
TABLE_ROWS_MAX = _DEFAULT["table_rows_max"]
REQ_RE = re.compile(r"(?:REQ[-_ ]?\d+|\bR\d{1,3}\b|요구\s*[:：]\s*\S|대응\s*요구)", re.IGNORECASE)
WARN = "[경고]"
NOTE = "[NOT INSPECTED]"
BLOCK_PREFIXES = ("[차단]",)


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return "sha256:" + h.hexdigest()


def shape_text(shape) -> str:
    """도형의 보이는 텍스트. 표(GraphicFrame)의 셀까지 읽는다.

    표는 has_text_frame이 False라 셀 텍스트가 밀도 계산에서 통째로 빠졌다 —
    조견표처럼 가장 빽빽한 장표가 '0자'로 집계돼 밀도 경고가 사실상 동작하지 않았다.
    """
    if shape.has_text_frame:
        return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
    if getattr(shape, "has_table", False) and shape.has_table:
        return "\n".join(cell.text for row in shape.table.rows for cell in row.cells).strip()
    return ""


def iter_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == 6 and hasattr(sh, "shapes"):  # GROUP
            yield from iter_shapes(sh.shapes)
        else:
            yield sh


def _group_transform(group) -> tuple[float, float, float, float] | None:
    """그룹의 (offset_x, offset_y, scale_x, scale_y). 못 읽으면 None.

    그룹 자식의 left/top은 자식 좌표계(chOff/chExt) 값이라 슬라이드 좌표가 아니다.
    변환 없이 비교하면 화면 안 도형을 '화면 밖'으로 잡거나 그 반대가 된다.
    """
    try:
        xfrm = group.element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr").find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
        a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        off, ext = xfrm.find(a + "off"), xfrm.find(a + "ext")
        ch_off, ch_ext = xfrm.find(a + "chOff"), xfrm.find(a + "chExt")
        cx, cy = int(ch_ext.get("cx")), int(ch_ext.get("cy"))
        if not cx or not cy:
            return None
        sx, sy = int(ext.get("cx")) / cx, int(ext.get("cy")) / cy
        return (int(off.get("x")) - int(ch_off.get("x")) * sx,
                int(off.get("y")) - int(ch_off.get("y")) * sy, sx, sy)
    except (AttributeError, TypeError, ValueError):
        return None


def placed_shapes(shapes, transform=(0.0, 0.0, 1.0, 1.0)):
    """(도형, 슬라이드 좌표 left/top/width/height). 그룹은 변환을 누적한다."""
    ox, oy, sx, sy = transform
    for sh in shapes:
        if sh.shape_type == 6 and hasattr(sh, "shapes"):  # GROUP
            inner = _group_transform(sh)
            if inner is None:
                # 변환을 못 읽으면 그룹 자체의 상자로만 판단한다(잘못된 좌표로
                # 자식을 재느니 검사 대상에서 뺀다).
                continue
            gx, gy, gsx, gsy = inner
            yield from placed_shapes(sh.shapes,
                                     (ox + gx * sx, oy + gy * sy, sx * gsx, sy * gsy))
            continue
        try:
            box = (ox + int(sh.left) * sx, oy + int(sh.top) * sy,
                   int(sh.width) * sx, int(sh.height) * sy)
        except (TypeError, ValueError):
            continue  # 좌표를 못 읽는 도형(플레이스홀더 상속 등)은 건너뛴다
        yield sh, box


def slide_kind(slide) -> str:
    names = {sh.name for sh in iter_shapes(slide.shapes)}
    if "COVER_BAND" in names and "SUBTITLE" in names:
        return "cover"
    if "COVER_BAND" in names:
        return "closing"
    if "SECTION_BAND" in names:
        return "section"
    if any(n.startswith("BODY_TOC") for n in names):
        return "toc"
    if "LEAD" in names or "TITLE" in names:
        return "body"
    return "unknown"


# 생성기가 넣는 상용구(위치·번호·캡션)만 제외한다. BODY_LEGEND는 구성도의 범례로
# 작성자가 쓴 내용이므로 폰트·밀도 검사 대상이다(예전엔 제외돼 본문 하한 미만으로 남았다).
META_SHAPES = {"HEADER", "FOOTER", "PAGENO", "CAPTION", "REQID", "STATUS"}
# 소형 텍스트(범례·간트 라벨)의 정의는 deck_profiles에 있다 — 생성기와 같은 정의를
# 읽어야 두 도구의 기준이 갈리지 않는다. 검사 대상이되 하한은 표 기준을 쓴다.


def _effective_pt(run, paragraph, shape) -> float | None:
    """실제로 적용되는 글자 크기. run에 없으면 문단, 그 다음 도형 기본값을 본다.

    run.font.size만 보면 문단 수준에서 크기를 준 텍스트가 통째로 '크기 미지정'으로
    빠져 폰트 하한 검사를 지나쳤다 — 축소한 장표가 그대로 통과하던 경로다.
    """
    for source in (run.font, paragraph.font,
                   getattr(getattr(shape, "text_frame", None), "_defRPr", None)):
        size = getattr(source, "size", None) if source is not None else None
        if size is not None:
            return size.pt
    return None


def min_font_pt(slide) -> tuple[float | None, float | None]:
    """(본문 최소 pt, 표 최소 pt). 캡션·헤더·푸터(상용구)는 제외.

    본문과 표는 허용 하한이 다르다(표가 더 작다). 하나로 합치면 표 하한이 본문에도
    적용돼, 본문을 표 크기까지 줄인 장표가 통과한다.
    """
    body: list[float] = []
    table: list[float] = []
    for sh in iter_shapes(slide.shapes):
        if sh.name in META_SHAPES:
            continue
        frames: list[tuple[object, list]] = []
        if sh.has_text_frame:
            frames.append((sh.text_frame, table if deck_profiles.is_small_text(sh.name) else body))
        if getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    frames.append((cell.text_frame, table))
        for tf, bucket in frames:
            for p in tf.paragraphs:
                for r in p.runs:
                    if not r.text.strip():
                        continue
                    pt = _effective_pt(r, p, sh)
                    if pt is not None:
                        bucket.append(pt)
    return (min(body) if body else None, min(table) if table else None)


def out_of_bounds(prs, slide) -> list[str]:
    """슬라이드 경계를 벗어난 텍스트 도형을 찾는다.

    폰트·글자 수 규칙만 보면 정상인데 상자가 화면 밖에 있어 아무것도 안 보이는 장표가
    통과했다. 렌더 PDF에도 안 나오므로 페이지 수 대조로도 걸리지 않는다.
    배경·장식(선·그림)은 의도적으로 걸칠 수 있으므로 텍스트가 있는 도형만 본다.
    """
    sw, sh_ = int(prs.slide_width), int(prs.slide_height)
    problems: list[str] = []
    for shape, (left, top, width, height) in placed_shapes(slide.shapes):
        if not shape_text(shape).strip():
            continue
        right, bottom = left + width, top + height
        if right <= 0 or bottom <= 0 or left >= sw or top >= sh_:
            problems.append(f"'{shape.name}' 완전히 화면 밖")
            continue
        over = max(0, -left) + max(0, right - sw)
        over_v = max(0, -top) + max(0, bottom - sh_)
        if width and over / width > 0.25:
            problems.append(f"'{shape.name}' 가로 {over / width:.0%} 잘림")
        elif height and over_v / height > 0.25:
            problems.append(f"'{shape.name}' 세로 {over_v / height:.0%} 잘림")
    return problems


def detect_profile(prs) -> tuple[str | None, str]:
    """(프로파일, 표시 상태). build_deck이 core properties에 남긴 표시를 읽는다.

    '표시 없음'과 '모르는 표시'를 구분한다. 둘을 같게 처리하면, 이 검사기가 모르는
    규격으로 만든 덱이 가장 느슨한 기본값으로 조용히 통과한다(fail-open).
    """
    try:
        raw = prs.core_properties.category
    except (AttributeError, ValueError):
        return None, "missing"
    return deck_profiles.read_stamp(raw)


def raw_stamp(prs) -> str:
    try:
        return str(prs.core_properties.category or "")
    except (AttributeError, ValueError):
        return ""


def stage_is_submission(stage: str) -> bool:
    return stage == "submission"


def lint(prs, *, max_pages: int | None, exclude_cover_toc: bool, min_font: float,
         require_req_ids: bool, stage: str, style: dict | None = None,
         min_table_font: float | None = None) -> list[str]:
    style = style or _DEFAULT
    if min_table_font is None:
        min_table_font = min(min_font, style["sizes"]["table"] - 1)
    lead_max = style["lead_max_chars"]
    density_max = style["density_max"]
    table_rows_max = style["table_rows_max"]
    items: list[str] = []
    slides = list(prs.slides)
    counted = 0
    for idx, slide in enumerate(slides, 1):
        kind = slide_kind(slide)
        # 도형을 한 번만 순회해 (이름, 텍스트) 쌍을 만든다. 이름 dict와 텍스트 list를
        # 따로 만들어 zip하면 같은 이름이 둘 이상일 때(간트 마일스톤 2개 이상) 짝이
        # 밀려 남의 텍스트를 세거나 본문을 통째로 누락했다.
        pairs = [(sh.name, shape_text(sh)) for sh in iter_shapes(slide.shapes)]
        names = {name: sh for name, sh in zip((n for n, _ in pairs), iter_shapes(slide.shapes))}
        texts = [t for _, t in pairs]
        all_text = "\n".join(t for t in texts if t)
        if not (exclude_cover_toc and kind in {"cover", "toc", "closing"}):
            counted += 1
        for problem in out_of_bounds(prs, slide):
            items.append(f"[차단] 슬라이드 {idx}: {problem} — 화면 밖 내용은 렌더에도 보이지 않는다")
        if not all_text.strip():
            pics = [sh for sh in iter_shapes(slide.shapes) if sh.shape_type == 13]
            items.append(f"{'[차단]' if not pics else WARN} 슬라이드 {idx}: "
                         f"{'빈 슬라이드' if not pics else '이미지 전용 슬라이드 — 텍스트 검사 불가, 렌더 육안 확인 필수'}")
            continue
        if kind in {"cover", "closing", "section", "toc"}:
            continue
        heuristic = not ({"TITLE", "LEAD"} <= names.keys())
        tag = " [휴리스틱]" if heuristic else ""
        if heuristic:
            # 외부 덱: 위치 순으로 상단 텍스트 2개를 제목·리드문으로 간주
            ordered = sorted((sh for sh in iter_shapes(slide.shapes) if shape_text(sh)),
                             key=lambda sh: (int(sh.top or 0), int(sh.left or 0)))
            title = shape_text(ordered[0]) if ordered else ""
            lead = shape_text(ordered[1]) if len(ordered) > 1 else ""
        else:
            title = shape_text(names["TITLE"])
            lead = shape_text(names["LEAD"])
        if not title:
            items.append(f"[차단] 슬라이드 {idx}: 제목 없음{tag}")
        if not lead:
            items.append(f"[차단] 슬라이드 {idx}: 리드문 없음 — 1페이지 1메시지 원칙{tag}")
        elif len(lead) > lead_max:
            items.append(f"{WARN} 슬라이드 {idx}: 리드문 {len(lead)}자 > {lead_max}자{tag}")
        elif lead.endswith(("설명합니다.", "설명합니다", "소개합니다.", "소개합니다")):
            items.append(f"{WARN} 슬라이드 {idx}: 리드문이 결론이 아닌 안내문('{lead[-6:]}')")
        if require_req_ids:
            req_text = shape_text(names["REQID"]) if "REQID" in names else all_text
            if not REQ_RE.search(req_text):
                level = "[차단]" if stage == "submission" else WARN
                items.append(f"{level} 슬라이드 {idx}: 대응 요구사항 ID(REQ-ID) 없음{tag}")
        # 밀도에서 제외하는 것과 폰트 검사에서 제외하는 것을 같은 정의(META_SHAPES)로
        # 맞춘다 — 캡션·REQ-ID는 생성기가 넣는 상용구이지 작성자가 쓴 본문이 아니다.
        body_chars = (sum(len(t) for n, t in pairs if n not in META_SHAPES)
                      if not heuristic else len(all_text))
        if body_chars > density_max:
            items.append(f"{WARN} 슬라이드 {idx}: 텍스트 {body_chars}자 > {density_max}자 — 도식·표로 압축 또는 분할")
        body_pt, table_pt = min_font_pt(slide)
        for measured, floor, role in ((body_pt, min_font, "본문"),
                                      (table_pt, min_table_font, "표")):
            if measured is not None and floor is not None and measured < floor:
                items.append(f"[차단] 슬라이드 {idx}: {role} 최소 폰트 {measured:g}pt < "
                             f"{floor:g}pt — 폰트 축소로 분량 회피 금지")
        for sh in iter_shapes(slide.shapes):
            if getattr(sh, "has_table", False) and sh.has_table:
                tbl = sh.table
                header = [c.text.strip() for c in tbl.rows[0].cells]
                if not any(header):
                    items.append(f"[차단] 슬라이드 {idx}: 표 헤더 행 비어 있음")
                if len(tbl.rows) - 1 > table_rows_max:
                    items.append(f"{WARN} 슬라이드 {idx}: 표 {len(tbl.rows) - 1}행 > {table_rows_max}행 — 분할·헤더 반복")
    if max_pages is not None:
        label = "본문" if exclude_cover_toc else "전체"
        if counted > max_pages:
            items.append(f"[차단] {label} {counted}장 > 페이지 제한 {max_pages}장")
        else:
            items.append(f"[정보] {label} {counted}장 / 제한 {max_pages}장")
    else:
        items.append(f"[정보] 전체 {len(slides)}장")
    return items


def render(pptx: Path, png_dir: Path | None) -> tuple[list[str], dict]:
    """soffice로 PDF 변환, 페이지 수 대조, 선택적으로 PNG 썸네일. (items, evidence)"""
    items: list[str] = []
    evidence: dict = {"tool": "", "pdf_pages": None, "png": []}
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        items.append(f"{NOTE} LibreOffice(soffice) 없음 — PDF 렌더 미수행. 사람이 PowerPoint에서 렌더 확인")
        return items, evidence
    with tempfile.TemporaryDirectory() as tmp:
        try:
            proc = subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(pptx)],
                                  capture_output=True, text=True, encoding="utf-8", errors="replace", timeout=240)
        except (subprocess.TimeoutExpired, OSError) as exc:
            items.append(f"[차단] 렌더 실패: {exc}")
            return items, evidence
        pdf = Path(tmp) / (pptx.stem + ".pdf")
        if proc.returncode != 0 or not pdf.is_file():
            items.append(f"[차단] 렌더 실패(soffice exit {proc.returncode}): {(proc.stderr or proc.stdout).strip()[:200]}")
            return items, evidence
        ver = subprocess.run([soffice, "--version"], capture_output=True, text=True,
                             encoding="utf-8", errors="replace").stdout.strip()
        evidence["tool"] = ver or "libreoffice"
        pages = None
        pdfinfo = shutil.which("pdfinfo")
        if pdfinfo:
            out = subprocess.run([pdfinfo, str(pdf)], capture_output=True, text=True,
                                 encoding="utf-8", errors="replace").stdout
            m = re.search(r"Pages:\s+(\d+)", out)
            pages = int(m.group(1)) if m else None
        if pages is None:
            pages = len(re.findall(rb"/Type\s*/Page[^s]", pdf.read_bytes()))
        evidence["pdf_pages"] = pages
        n_slides = len(Presentation(str(pptx)).slides)
        if pages != n_slides:
            items.append(f"[차단] 렌더 페이지 {pages} ≠ 슬라이드 {n_slides} — 숨김 슬라이드·렌더 누락 확인")
        else:
            items.append(f"[정보] 렌더 OK: PDF {pages}p = 슬라이드 {n_slides}장 ({evidence['tool']})")
        if png_dir:
            png_dir.mkdir(parents=True, exist_ok=True)
            pdftoppm = shutil.which("pdftoppm")
            if pdftoppm:
                subprocess.run([pdftoppm, "-r", "60", "-png", str(pdf), str(png_dir / "slide")],
                               capture_output=True, timeout=240)
                evidence["png"] = sorted(p.name for p in png_dir.glob("slide-*.png"))
                items.append(f"[정보] 썸네일 {len(evidence['png'])}장 → {png_dir} (육안 검사: 잘림·겹침·폰트 대체)")
            else:
                shutil.copy(pdf, png_dir / pdf.name)
                items.append(f"{NOTE} pdftoppm 없음 — PDF만 저장({png_dir / pdf.name}). 육안 검사 필요")
    return items, evidence


def is_blocking(item: str) -> bool:
    return item.startswith(BLOCK_PREFIXES)


def main(argv: list[str] | None = None) -> int:
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8", errors="replace")
        except (AttributeError, ValueError):
            pass
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--max-pages", type=int, help="페이지 제한(RFP)")
    ap.add_argument("--exclude-cover-toc", action="store_true", help="표지·목차·마무리를 제한 산정에서 제외")
    ap.add_argument("--min-font", type=float,
                    help="본문 최소 폰트(pt). 기본값은 프로파일에서 유도한다. 캡션·헤더·푸터는 제외")
    ap.add_argument("--profile", choices=sorted(deck_profiles.PROFILES),
                    help="시각 규격을 명시한다. 생략하면 PPTX에 남은 표시를 읽고, "
                         f"표시가 없으면 {deck_profiles.DEFAULT_PROFILE}")
    ap.add_argument("--require-req-ids", action="store_true", help="본문 장표 REQ-ID 표기 필수(공공·금융 A/B)")
    ap.add_argument("--stage", choices=["draft", "submission"], default="submission")
    ap.add_argument("--render", action="store_true", help="LibreOffice로 PDF 렌더·페이지 대조")
    ap.add_argument("--png-dir", type=Path, help="렌더 썸네일 저장 폴더(--render 필요)")
    ap.add_argument("--emit-render", type=Path, help="audit용 render 블록 JSON 출력 경로")
    a = ap.parse_args(argv)
    if not a.pptx.is_file():
        print(f"파일 없음: {a.pptx}", file=sys.stderr)
        return 2
    try:
        prs = Presentation(str(a.pptx))
    except Exception as exc:  # pptx 패키지 손상 등
        print(f"검사 불가(파일 손상·형식): {exc}", file=sys.stderr)
        return 2
    stamped, stamp_state = detect_profile(prs)
    profile = a.profile or stamped or deck_profiles.DEFAULT_PROFILE
    style = deck_profiles.get(profile)
    # 최소 폰트는 프로파일에서 유도한다 — 상수로 고정하면 발표본(18pt)을 상세본
    # 기준(9pt)으로 재거나, 그 반대로 정상 산출물을 차단한다.
    min_font = a.min_font if a.min_font is not None else deck_profiles.min_body_font(profile)
    min_table_font = (min(a.min_font, deck_profiles.min_table_font(profile))
                      if a.min_font is not None else deck_profiles.min_table_font(profile))
    # 실제 로더로 한 번 열어본다 — 구조 검사가 놓치는 관계·스키마 결함은 여기서 걸린다.
    load_problem = quality_gate.load_check(a.pptx)
    if load_problem:
        print(f"검사 불가: {load_problem}", file=sys.stderr)
        return 2
    source = "지정" if a.profile else ("파일 표시" if stamped else "기본값 추정")
    items = [f"[정보] 규격: {style['label']}({profile}, {source}) · 최소 폰트 {min_font}pt · "
             f"표 {min_table_font}pt · 밀도 {style['density_max']}자"]
    if a.profile and stamped and a.profile != stamped:
        items.append(f"{WARN} 지정한 규격({a.profile})이 파일에 남은 표시({stamped})와 다르다 — "
                     "다른 규격으로 만든 파일을 검사하고 있는지 확인한다")
    if not a.profile and stamp_state != "known":
        # 표시가 없거나 모르는 값이면 가장 느슨한 기본값으로 재게 된다. 제출 단계에서는
        # 추정으로 통과시키지 않는다(미검사 ≠ 통과).
        detail = ("규격 표시가 없다" if stamp_state == "missing"
                  else f"이 검사기가 모르는 규격 표시({raw_stamp(prs)})다")
        level = "[차단]" if stage_is_submission(a.stage) else WARN
        items.append(f"{level} {detail} — {deck_profiles.DEFAULT_PROFILE} 기준으로 재고 있다. "
                     "`--profile`로 규격을 명시하거나 build_deck으로 다시 생성한다")
    items += lint(prs, max_pages=a.max_pages, exclude_cover_toc=a.exclude_cover_toc,
                  min_font=min_font, min_table_font=min_table_font,
                  require_req_ids=a.require_req_ids, stage=a.stage, style=style)
    evidence: dict = {}
    rendered = False
    if a.render or a.png_dir or a.emit_render:
        r_items, evidence = render(a.pptx, a.png_dir)
        items += r_items
        rendered = bool(evidence.get("tool")) and not any(is_blocking(i) for i in r_items)
    blockers = [i for i in items if is_blocking(i)]
    for i in items:
        print(i)
    if a.emit_render:
        block = {
            "verified": rendered and not blockers,
            # 렌더 성공(기계)과 육안 승인(사람)은 다른 사실이다. PDF 변환이 됐다는 것이
            # 장표를 사람이 보고 승인했다는 뜻은 아니므로 따로 기록한다.
            "render_succeeded": rendered,
            "layout_checked": True,
            "visual_review_approved": False,
            "visual_reviewer": "",
            "output_profile": profile,
            "artifact_hash": sha256(a.pptx),
            "tool": f"deck_check.py + {evidence.get('tool') or 'no-renderer'}",
            "evidence": [i for i in items if not i.startswith(WARN)] + (
                [f"png:{n}" for n in evidence.get("png", [])[:3]] if evidence.get("png") else []),
        }
        if not rendered:
            block["evidence"].append("NOT INSPECTED: PDF 렌더 미수행 — verified=false")
        block["evidence"].append(
            "육안 검토 미완료: 썸네일을 확인한 뒤 visual_review_approved=true와 "
            "visual_reviewer를 사람이 기록한다")
        a.emit_render.write_text(json.dumps(block, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        print(f"[정보] render 블록 → {a.emit_render} (verified={block['verified']})")
    if blockers:
        print(f"차단 — {len(blockers)}건")
        return 1
    print("통과 (레이아웃 린트" + (" + 렌더" if rendered else "") + " — 도식 의미·설득력은 사람이 검토)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
