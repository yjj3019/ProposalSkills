#!/usr/bin/env python3
"""slides.json → 제안서 PPTX 생성기 (visual-style.md 규격의 결정론적 구현).

사용법:
    python3 build_deck.py slides.json -o 제안서.pptx [--template 사내양식.pptx] [--strict]

매 실행마다 레이아웃·색·폰트를 재발명하지 않도록, 장표 구성(헤더 breadcrumb → 제목 →
리드문 → 본문 → 캡션/출처/REQ-ID → 푸터)을 고정 좌표·고정 토큰으로 그린다.
모든 도형에 역할 이름(HEADER/TITLE/LEAD/BODY/CAPTION/REQID/FOOTER)을 부여해
deck_check.py가 리드문·REQ-ID 존재를 기계 검사할 수 있게 한다.

slides.json 스키마 요약(상세: references/deck-production.md):
  {"meta": {"title","subtitle","buyer","bidder","date","doc_name","page_limit","palette":{...},"font"},
   "slides": [ {"type": "cover|toc|section|matrix|table|process|zones|gantt|staff|cards|bullets|closing",
                "breadcrumb","title","lead","caption","source","req_ids":[...],"notes","status_tag", ...} ]}

종료 코드: 0=생성, 1=--strict 위반(리드문 없는 본문 장표 등), 2=입력 오류.
"""
from __future__ import annotations

import argparse
import json
import math
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Emu, Pt
except ImportError:  # pragma: no cover
    print("python-pptx가 필요합니다: pip install python-pptx", file=sys.stderr)
    raise SystemExit(2)

# ---- 토큰 (visual-style.md §0 폴백 토큰과 동일) --------------------------------
DEFAULT_PALETTE = {
    "primary": "1F3864", "tint1": "8FAADC", "tint2": "D6E0F0", "tint3": "EDF1F8",
    "text": "202020", "muted": "595959", "warn": "C00000", "ok": "1F7A3D", "white": "FFFFFF",
}
DEFAULT_FONT = "맑은 고딕"
SIZE = {"title": 22, "lead": 13, "body": 11, "table": 10, "caption": 8.5, "header": 9,
        "footer": 9, "cover_title": 32, "cover_sub": 16, "section": 28}
LEAD_MAX_CHARS = 60          # writing-style.md: 리드문 60자 이내
MATRIX_ROWS_PER_SLIDE = 12   # 조견표 자동 분할 기준

# 16:9 (13.333in × 7.5in) 고정 그리드 — 좌표는 EMU
IN = 914400
W, H = int(13.333 * IN), int(7.5 * IN)
M = int(0.5 * IN)                     # 좌우 여백
HEADER_Y, HEADER_H = int(0.25 * IN), int(0.3 * IN)
TITLE_Y, TITLE_H = int(0.6 * IN), int(0.55 * IN)
LEAD_Y, LEAD_H = int(1.2 * IN), int(0.65 * IN)
BODY_Y = int(1.95 * IN)
CAPTION_Y, CAPTION_H = int(6.55 * IN), int(0.35 * IN)
FOOTER_Y, FOOTER_H = int(6.95 * IN), int(0.3 * IN)
BODY_H = CAPTION_Y - BODY_Y - int(0.1 * IN)
BODY_W = W - 2 * M

BODY_TYPES = {"matrix", "table", "process", "zones", "gantt", "staff", "cards", "bullets"}
ALL_TYPES = BODY_TYPES | {"cover", "toc", "section", "closing"}


def rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6.lstrip("#").upper())


class DeckBuilder:
    def __init__(self, spec: dict, template: Path | None, strict: bool):
        self.spec = spec
        self.meta = spec.get("meta", {})
        self.pal = {**DEFAULT_PALETTE, **{k: str(v).lstrip("#").upper()
                                          for k, v in (self.meta.get("palette") or {}).items()}}
        self.font = self.meta.get("font") or DEFAULT_FONT
        self.strict = strict
        self.violations: list[str] = []
        self.warnings: list[str] = []
        self.page = 0
        if template:
            self.prs = Presentation(str(template))
            self.warnings.append(f"사용자 템플릿 사용: {template.name} — 마스터 배경·로고는 템플릿 것을 따른다")
            if abs(int(self.prs.slide_width) - W) > IN // 10 or abs(int(self.prs.slide_height) - H) > IN // 10:
                raise ValueError("템플릿은 16:9(13.333in×7.5in)여야 한다 — 다른 비율은 좌표 그리드와 맞지 않는다")
        else:
            self.prs = Presentation()
            self.prs.slide_width, self.prs.slide_height = Emu(W), Emu(H)
        self.blank = self._blank_layout()

    # ---- 기본 도형 헬퍼 ------------------------------------------------------
    def _blank_layout(self):
        for layout in self.prs.slide_layouts:
            if layout.name.lower() in {"blank", "빈 화면", "빈화면"}:
                return layout
        return self.prs.slide_layouts[len(self.prs.slide_layouts) - 1]

    def _text(self, slide, name, x, y, w, h, text, size, *, bold=False, color=None,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, wrap=True):
        box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        box.name = name
        if fill:
            box.fill.solid()
            box.fill.fore_color.rgb = rgb(fill)
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.auto_size = MSO_AUTO_SIZE.NONE  # 렌더러 자동 축소로 상자 높이가 변하지 않게 고정
        tf.vertical_anchor = anchor
        for attr in ("margin_left", "margin_right"):
            setattr(tf, attr, Emu(int(0.05 * IN)))
        tf.margin_top = tf.margin_bottom = Emu(int(0.03 * IN))
        lines = text if isinstance(text, list) else [text]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = str(line)
            r.font.name = self.font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = rgb(color or self.pal["text"])
        return box

    def _rect(self, slide, name, x, y, w, h, fill, *, line=None, shape=MSO_SHAPE.RECTANGLE):
        s = slide.shapes.add_shape(shape, Emu(x), Emu(y), Emu(w), Emu(h))
        s.name = name
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
        if line:
            s.line.color.rgb = rgb(line)
            s.line.width = Pt(0.75)
        else:
            s.line.fill.background()
        s.shadow.inherit = False
        return s

    def _shape_text(self, shape, text, size, *, bold=False, color=None, align=PP_ALIGN.CENTER):
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lines = text if isinstance(text, list) else [text]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = str(line)
            r.font.name = self.font
            r.font.size = Pt(size)
            r.font.bold = bold if i == 0 else False
            r.font.color.rgb = rgb(color or self.pal["text"])

    def _table(self, slide, name, x, y, w, h, columns, rows, *, col_widths=None,
               size=SIZE["table"], align_right_cols=()):
        n_rows, n_cols = len(rows) + 1, len(columns)
        gt = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(w), Emu(h))
        gt.name = name
        table = gt.table
        if col_widths:
            total = sum(col_widths)
            for i, cw in enumerate(col_widths):
                table.columns[i].width = Emu(int(w * cw / total))
        row_h = int(h / n_rows)
        for r in range(n_rows):
            table.rows[r].height = Emu(row_h)
        for c, col in enumerate(columns):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(self.pal["primary"])
            self._cell(cell, col, size, bold=True, color=self.pal["white"], align=PP_ALIGN.CENTER)
        for r, row in enumerate(rows, 1):
            for c in range(n_cols):
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(self.pal["tint3"] if r % 2 == 0 else self.pal["white"])
                val = row[c] if c < len(row) else ""
                align = PP_ALIGN.RIGHT if c in align_right_cols else PP_ALIGN.LEFT
                self._cell(cell, val, size, align=align)
        return gt

    def _cell(self, cell, text, size, *, bold=False, color=None, align=PP_ALIGN.LEFT):
        cell.margin_left = cell.margin_right = Emu(int(0.05 * IN))
        cell.margin_top = cell.margin_bottom = Emu(int(0.02 * IN))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = "" if text is None else str(text)
        r.font.name = self.font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color or self.pal["text"])

    # ---- 공통 프레임 ----------------------------------------------------------
    def _frame(self, s: dict):
        """헤더 → 제목 → 리드문 → (본문은 호출자) → 캡션/REQ-ID → 푸터."""
        slide = self.prs.slides.add_slide(self.blank)
        self.page += 1
        doc = self.meta.get("doc_name") or self.meta.get("title") or ""
        crumb = s.get("breadcrumb", "")
        self._text(slide, "HEADER", M, HEADER_Y, BODY_W - int(2.2 * IN), HEADER_H,
                   f"{doc}  |  {crumb}" if crumb else doc, SIZE["header"], color=self.pal["muted"])
        if s.get("status_tag"):  # 유형 C: 우상단 항목 번호·상태
            self._text(slide, "STATUS", W - M - int(2.2 * IN), HEADER_Y, int(2.2 * IN), HEADER_H,
                       s["status_tag"], SIZE["header"], bold=True, color=self.pal["primary"],
                       align=PP_ALIGN.RIGHT)
        self._text(slide, "TITLE", M, TITLE_Y, BODY_W, TITLE_H, s.get("title", ""), SIZE["title"],
                   bold=True, color=self.pal["primary"], anchor=MSO_ANCHOR.MIDDLE)
        line = self._rect(slide, "TITLE_RULE", M, TITLE_Y + TITLE_H, BODY_W, int(0.02 * IN), self.pal["primary"])
        lead = s.get("lead", "")
        if not lead and s.get("type") in BODY_TYPES:
            self.violations.append(f"슬라이드 {self.page} '{s.get('title', '')}': 리드문 없음")
        if lead and len(lead) > LEAD_MAX_CHARS:
            self.warnings.append(f"슬라이드 {self.page}: 리드문 {len(lead)}자 (권장 ≤{LEAD_MAX_CHARS})")
        self._text(slide, "LEAD", M, LEAD_Y, BODY_W, LEAD_H, lead, SIZE["lead"], bold=True,
                   anchor=MSO_ANCHOR.MIDDLE)
        cap_parts = []
        if s.get("caption"):
            cap_parts.append(s["caption"])
        if s.get("source"):
            cap_parts.append(f"출처: {s['source']}")
        self._text(slide, "CAPTION", M, CAPTION_Y, BODY_W - int(3.2 * IN), CAPTION_H,
                   "  ·  ".join(cap_parts), SIZE["caption"], color=self.pal["muted"])
        req = s.get("req_ids") or []
        self._text(slide, "REQID", W - M - int(3.2 * IN), CAPTION_Y, int(3.2 * IN), CAPTION_H,
                   ("대응 요구: " + ", ".join(map(str, req))) if req else "", SIZE["caption"],
                   color=self.pal["muted"], align=PP_ALIGN.RIGHT)
        if not req and s.get("type") in BODY_TYPES and self.meta.get("require_req_ids", True):
            self.warnings.append(f"슬라이드 {self.page} '{s.get('title', '')}': REQ-ID 없음 (공공·금융 A/B 유형은 필수)")
        bidder = self.meta.get("bidder", "")
        self._text(slide, "FOOTER", M, FOOTER_Y, BODY_W - int(1.0 * IN), FOOTER_H, bidder,
                   SIZE["footer"], color=self.pal["muted"])
        self._text(slide, "PAGENO", W - M - int(1.0 * IN), FOOTER_Y, int(1.0 * IN), FOOTER_H,
                   str(self.page), SIZE["footer"], color=self.pal["muted"], align=PP_ALIGN.RIGHT)
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(s["notes"])
        return slide

    # ---- 슬라이드 유형 ---------------------------------------------------------
    def cover(self, s):
        slide = self.prs.slides.add_slide(self.blank)
        self.page += 1
        self._rect(slide, "COVER_BAND", 0, int(2.3 * IN), W, int(2.6 * IN), self.pal["primary"])
        self._text(slide, "TITLE", M, int(2.5 * IN), BODY_W, int(1.3 * IN),
                   s.get("title") or self.meta.get("title", ""), SIZE["cover_title"], bold=True,
                   color=self.pal["white"], anchor=MSO_ANCHOR.MIDDLE)
        self._text(slide, "SUBTITLE", M, int(3.8 * IN), BODY_W, int(0.9 * IN),
                   s.get("subtitle") or self.meta.get("subtitle", ""), SIZE["cover_sub"],
                   color=self.pal["white"], anchor=MSO_ANCHOR.MIDDLE)
        lines = [f"제출: {self.meta.get('buyer', '')}", f"제안: {self.meta.get('bidder', '')}",
                 self.meta.get("date", "")]
        self._text(slide, "COVER_META", M, int(5.4 * IN), BODY_W, int(1.2 * IN),
                   [line for line in lines if line.strip(": ")], SIZE["body"], color=self.pal["muted"])
        return slide

    def toc(self, s):
        slide = self._frame({**s, "type": "toc", "title": s.get("title", "목차"), "lead": s.get("lead", "")})
        items = s.get("items", [])
        half = math.ceil(len(items) / 2)
        for col, chunk in enumerate((items[:half], items[half:])):
            if not chunk:
                continue
            x = M + col * (BODY_W // 2 + int(0.2 * IN))
            self._text(slide, f"BODY_TOC{col}", x, BODY_Y, BODY_W // 2 - int(0.2 * IN), BODY_H,
                       [str(i) for i in chunk], SIZE["body"] + 1)
        return slide

    def section(self, s):
        slide = self.prs.slides.add_slide(self.blank)
        self.page += 1
        self._rect(slide, "SECTION_BAND", 0, 0, int(0.35 * IN), H, self.pal["primary"])
        self._text(slide, "SECTION_NO", M + int(0.3 * IN), int(2.4 * IN), BODY_W, int(0.8 * IN),
                   s.get("no", ""), SIZE["section"], bold=True, color=self.pal["tint1"])
        self._text(slide, "TITLE", M + int(0.3 * IN), int(3.2 * IN), BODY_W, int(1.0 * IN),
                   s.get("title", ""), SIZE["section"], bold=True, color=self.pal["primary"])
        if s.get("items"):
            self._text(slide, "SECTION_ITEMS", M + int(0.3 * IN), int(4.3 * IN), BODY_W, int(2.2 * IN),
                       [f"{i}" for i in s["items"]], SIZE["body"], color=self.pal["muted"])
        return slide

    def table(self, s):
        slide = self._frame(s)
        columns, rows = s.get("columns", []), s.get("rows", [])
        if not columns or not rows:
            self.violations.append(f"슬라이드 {self.page}: table에 columns/rows 없음")
            return slide
        h = min(BODY_H, int((len(rows) + 1) * 0.38 * IN))
        self._table(slide, "BODY_TABLE", M, BODY_Y, BODY_W, h, columns, rows,
                    col_widths=s.get("col_widths"), align_right_cols=tuple(s.get("right_cols", [])))
        if len(rows) > 14:
            self.warnings.append(f"슬라이드 {self.page}: 표 {len(rows)}행 — 분할 권장(visual-style §4)")
        return slide

    def matrix(self, s):
        """조견표: 요구사항 대응표. rows가 많으면 자동 분할, 헤더 반복, '(계속)' 표기."""
        rows = s.get("rows", [])
        columns = s.get("columns") or ["ID", "제안요건", "수용여부", "제안내용/응답위치", "비고"]
        keys = s.get("keys") or ["id", "text", "support", "response_loc", "note"]
        per = int(s.get("rows_per_slide") or MATRIX_ROWS_PER_SLIDE)
        chunks = [rows[i:i + per] for i in range(0, len(rows), per)] or [[]]
        for i, chunk in enumerate(chunks):
            title = s.get("title", "요구사항 대응표(조견표)") + (f" (계속 {i + 1}/{len(chunks)})" if len(chunks) > 1 and i else
                                                          (f" (1/{len(chunks)})" if len(chunks) > 1 else ""))
            sub = {**s, "title": title, "req_ids": s.get("req_ids") or [r.get("id", "") for r in chunk]}
            slide = self._frame(sub)
            data = [[r.get(k, "") for k in keys] for r in chunk]
            h = min(BODY_H, int((len(data) + 1) * 0.36 * IN))
            self._table(slide, "BODY_MATRIX", M, BODY_Y, BODY_W, h, columns, data,
                        col_widths=s.get("col_widths") or [1, 5, 1.3, 4, 2.5])
        return None

    def process(self, s):
        """단계 프로세스: 가로 셰브론/박스 + 단계 설명."""
        slide = self._frame(s)
        steps = s.get("steps", [])
        if not steps:
            self.violations.append(f"슬라이드 {self.page}: process에 steps 없음")
            return slide
        n = len(steps)
        gap = int(0.12 * IN)
        bw = int((BODY_W - gap * (n - 1)) / n)
        top_h = int(0.9 * IN)
        for i, st in enumerate(steps):
            x = M + i * (bw + gap)
            box = self._rect(slide, f"BODY_STEP{i + 1}", x, BODY_Y, bw, top_h,
                             self.pal["primary"] if i == 0 else self.pal["tint1"], shape=MSO_SHAPE.CHEVRON if n <= 6 else MSO_SHAPE.RECTANGLE)
            self._shape_text(box, [f"{i + 1}. {st.get('title', '')}"], SIZE["body"], bold=True,
                             color=self.pal["white"] if i == 0 else self.pal["text"])
            desc = st.get("desc", [])
            desc_h = min(BODY_H - top_h - int(0.15 * IN), int(2.6 * IN))
            self._text(slide, f"BODY_STEPDESC{i + 1}", x, BODY_Y + top_h + int(0.15 * IN), bw, desc_h,
                       [f"• {d}" for d in (desc if isinstance(desc, list) else [desc])], SIZE["body"] - 1,
                       fill=self.pal["tint3"])
        return slide

    def zones(self, s):
        """구성도(아키텍처): 영역(zone) 박스 세로 배열 + 영역 내 구성요소 카드. 범례 포함."""
        slide = self._frame(s)
        zones = s.get("zones", [])
        if not zones:
            self.violations.append(f"슬라이드 {self.page}: zones 없음")
            return slide
        legend_h = int(0.35 * IN) if s.get("legend") else 0
        gap = int(0.1 * IN)
        zh = int((BODY_H - legend_h - gap * (len(zones) - 1)) / len(zones))
        label_w = int(1.6 * IN)
        for zi, z in enumerate(zones):
            y = BODY_Y + zi * (zh + gap)
            self._rect(slide, f"BODY_ZONE{zi + 1}", M, y, BODY_W, zh, self.pal["tint3"], line=self.pal["tint1"])
            lab = self._rect(slide, f"BODY_ZONELABEL{zi + 1}", M, y, label_w, zh, self.pal["primary"])
            self._shape_text(lab, [z.get("title", "")], SIZE["body"], bold=True, color=self.pal["white"])
            items = z.get("items", [])
            if items:
                iw = int((BODY_W - label_w - gap * (len(items) + 1)) / len(items))
                ih = zh - 2 * gap
                for ii, it in enumerate(items):
                    x = M + label_w + gap + ii * (iw + gap)
                    card = self._rect(slide, f"BODY_Z{zi + 1}I{ii + 1}", x, y + gap, iw, ih, self.pal["white"],
                                      line=self.pal["tint1"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
                    if isinstance(it, dict):
                        self._shape_text(card, [it.get("title", "")] + [str(d) for d in it.get("desc", [])],
                                         SIZE["body"] - 1, bold=True)
                    else:
                        self._shape_text(card, [str(it)], SIZE["body"] - 1, bold=True)
        if s.get("legend"):
            self._text(slide, "BODY_LEGEND", M, BODY_Y + BODY_H - legend_h, BODY_W, legend_h,
                       "범례: " + "  ·  ".join(s["legend"]), SIZE["caption"], color=self.pal["muted"])
        return slide

    def gantt(self, s):
        """추진일정 간트: 월 단위 컬럼, 단계 행, 마일스톤 ◆."""
        slide = self._frame(s)
        months = int(s.get("months") or 0)
        tasks = s.get("tasks", [])
        if not months or not tasks:
            self.violations.append(f"슬라이드 {self.page}: gantt에 months/tasks 없음")
            return slide
        label_w = int(2.6 * IN)
        grid_x = M + label_w
        grid_w = BODY_W - label_w
        head_h = int(0.35 * IN)
        row_h = min(int(0.42 * IN), int((BODY_H - head_h) / len(tasks)))
        mw = grid_w / months
        labels = s.get("month_labels") or [f"M{i + 1}" for i in range(months)]
        for mi in range(months):
            x = int(grid_x + mi * mw)
            hdr = self._rect(slide, f"BODY_GM{mi + 1}", x, BODY_Y, int(mw), head_h,
                             self.pal["primary"], line=self.pal["white"])
            self._shape_text(hdr, [labels[mi] if mi < len(labels) else f"M{mi + 1}"], SIZE["table"] - 1,
                             bold=True, color=self.pal["white"])
        for ti, t in enumerate(tasks):
            y = BODY_Y + head_h + ti * row_h
            self._rect(slide, f"BODY_GROW{ti + 1}", M, y, BODY_W, row_h,
                       self.pal["tint3"] if ti % 2 else self.pal["white"], line=self.pal["tint2"])
            self._text(slide, f"BODY_GLABEL{ti + 1}", M, y, label_w, row_h, t.get("name", ""),
                       SIZE["table"], bold=bool(t.get("phase")), anchor=MSO_ANCHOR.MIDDLE)
            start, end = float(t.get("start", 1)), float(t.get("end", t.get("start", 1)))
            bx = int(grid_x + (start - 1) * mw)
            bw = max(int((end - start + 1) * mw) - int(0.04 * IN), int(0.08 * IN))
            bar = self._rect(slide, f"BODY_GBAR{ti + 1}", bx, y + int(row_h * 0.25), bw, int(row_h * 0.5),
                             self.pal["primary"] if t.get("phase") else self.pal["tint1"])
            if t.get("label"):
                self._shape_text(bar, [t["label"]], SIZE["table"] - 1, color=self.pal["white"] if t.get("phase") else self.pal["text"])
            for ms in t.get("milestones", []):
                mx = int(grid_x + (float(ms.get("at", end)) - 0.5) * mw)
                dia = self._rect(slide, f"BODY_GMS{ti + 1}", mx - int(0.09 * IN), y + int(row_h * 0.2),
                                 int(0.18 * IN), int(row_h * 0.6), self.pal["warn"], shape=MSO_SHAPE.DIAMOND)
                if ms.get("label"):
                    self._text(slide, f"BODY_GMSL{ti + 1}", mx + int(0.1 * IN), y, int(1.6 * IN), row_h,
                               ms["label"], SIZE["table"] - 1, color=self.pal["warn"], anchor=MSO_ANCHOR.MIDDLE)
        return slide

    def staff(self, s):
        """인력 프로필 표: 성명 | 역할 | 등급 | 경력 | 자격 | 투입(M/M)."""
        people = s.get("people", [])
        columns = s.get("columns") or ["성명", "역할", "등급", "경력", "보유 자격", "투입(M/M)"]
        keys = s.get("keys") or ["name", "role", "grade", "years", "certs", "mm"]
        rows = [[(", ".join(p[k]) if isinstance(p.get(k), list) else p.get(k, "")) for k in keys] for p in people]
        return self.table({**s, "columns": columns, "rows": rows, "col_widths": s.get("col_widths") or [1.2, 2.2, 1, 1, 3.5, 1.1],
                           "right_cols": [len(columns) - 1]})

    def cards(self, s):
        """차별점·기대효과 카드(2~4열): 제목 + 값/한 줄 + 근거."""
        slide = self._frame(s)
        items = s.get("items", [])
        if not items:
            self.violations.append(f"슬라이드 {self.page}: cards에 items 없음")
            return slide
        n = len(items)
        cols = min(n, 4)
        rows_n = math.ceil(n / cols)
        gap = int(0.15 * IN)
        cw = int((BODY_W - gap * (cols - 1)) / cols)
        ch = min(int((BODY_H - gap * (rows_n - 1)) / rows_n), int(3.4 * IN))
        for i, it in enumerate(items):
            r, c = divmod(i, cols)
            x, y = M + c * (cw + gap), BODY_Y + r * (ch + gap)
            self._rect(slide, f"BODY_CARD{i + 1}", x, y, cw, ch, self.pal["tint3"], line=self.pal["tint1"])
            hdr = self._rect(slide, f"BODY_CARDHDR{i + 1}", x, y, cw, int(0.45 * IN), self.pal["primary"])
            self._shape_text(hdr, [it.get("title", "")], SIZE["body"], bold=True, color=self.pal["white"])
            body_lines = []
            if it.get("value"):
                body_lines.append(str(it["value"]))
            for d in it.get("desc", []) if isinstance(it.get("desc"), list) else [it.get("desc", "")]:
                if d:
                    body_lines.append(f"• {d}")
            if it.get("evidence"):
                body_lines.append(f"근거: {it['evidence']}")
            self._text(slide, f"BODY_CARDTXT{i + 1}", x, y + int(0.5 * IN), cw, ch - int(0.55 * IN), body_lines,
                       SIZE["body"] - 1)
        return slide

    def bullets(self, s):
        slide = self._frame(s)
        items = s.get("items", [])
        self.warnings.append(f"슬라이드 {self.page}: bullets 유형 — 텍스트 나열은 최소화(도식·표 권장)")
        self._text(slide, "BODY_BULLETS", M, BODY_Y, BODY_W, BODY_H, [f"• {i}" for i in items], SIZE["body"] + 1)
        if sum(len(str(i)) for i in items) > 450:
            self.warnings.append(f"슬라이드 {self.page}: 텍스트 과밀({sum(len(str(i)) for i in items)}자) — 분할 권장")
        return slide

    def closing(self, s):
        slide = self.prs.slides.add_slide(self.blank)
        self.page += 1
        self._rect(slide, "COVER_BAND", 0, int(2.8 * IN), W, int(1.8 * IN), self.pal["primary"])
        self._text(slide, "TITLE", M, int(3.0 * IN), BODY_W, int(1.4 * IN), s.get("title", "감사합니다"),
                   SIZE["cover_title"], bold=True, color=self.pal["white"], anchor=MSO_ANCHOR.MIDDLE,
                   align=PP_ALIGN.CENTER)
        if s.get("lead"):
            self._text(slide, "LEAD", M, int(4.8 * IN), BODY_W, int(0.8 * IN), s["lead"], SIZE["lead"],
                       align=PP_ALIGN.CENTER, color=self.pal["muted"])
        return slide

    # ---- 실행 -----------------------------------------------------------------
    def build(self) -> None:
        slides = self.spec.get("slides")
        if not isinstance(slides, list) or not slides:
            raise ValueError("slides must be a non-empty array")
        for i, s in enumerate(slides, 1):
            if not isinstance(s, dict):
                raise ValueError(f"slide {i} must be an object")
            t = s.get("type")
            if t not in ALL_TYPES:
                raise ValueError(f"slide {i}: unsupported type {t!r} (allowed: {sorted(ALL_TYPES)})")
            getattr(self, t)(s)
        limit = self.meta.get("page_limit")
        if limit and self.page > int(limit):
            self.violations.append(f"총 {self.page}장 > 페이지 제한 {limit}장 (표지·간지 포함 여부는 RFP 규정 확인)")

    def save(self, out: Path) -> None:
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))


def main(argv: list[str] | None = None) -> int:
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8", errors="replace")
        except (AttributeError, ValueError):
            pass
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("spec", type=Path, help="slides.json")
    ap.add_argument("-o", "--output", type=Path, required=True, help="출력 PPTX 경로")
    ap.add_argument("--template", type=Path, help="사내 양식 PPTX(마스터·배경 사용, 빈 레이아웃 필요)")
    ap.add_argument("--strict", action="store_true", help="리드문 누락 등 위반 시 exit 1")
    a = ap.parse_args(argv)
    try:
        spec = json.loads(a.spec.read_text(encoding="utf-8-sig"))
        if a.template and not a.template.is_file():
            raise ValueError(f"템플릿 없음: {a.template}")
        b = DeckBuilder(spec, a.template, a.strict)
        b.build()
        b.save(a.output)
    except (OSError, ValueError) as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 2
    print(f"wrote {a.output} ({b.page} slides)")
    for w in b.warnings:
        print(f"  경고: {w}")
    for v in b.violations:
        print(f"  위반: {v}")
    if b.violations and a.strict:
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
