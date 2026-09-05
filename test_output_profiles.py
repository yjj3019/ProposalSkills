"""회귀 테스트 — 산출물 종류별 시각 규격(F08).

읽는 조건이 다르면 규격이 달라야 한다. 동시에 생성기와 검사기가 같은 정의를 읽어야
"생성기는 통과, 검사기는 차단" 같은 어긋남이 생기지 않는다(실제로 표 행 14/15,
밀도 450/600으로 갈려 있었다).
"""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

REPO = Path(__file__).resolve().parent
DOC = REPO / "skills/create-proposal-document"
BD = DOC / "scripts/build_deck.py"
DC = DOC / "scripts/deck_check.py"
FIX = DOC / "fixtures/e2e-mini-rfp/slides.json"

sys.path.insert(0, str(DOC / "scripts"))
import deck_profiles as dp  # noqa: E402

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:  # pragma: no cover
    HAS_PPTX = False


def run(*args: object) -> subprocess.CompletedProcess:
    return subprocess.run([sys.executable, *map(str, args)], capture_output=True, text=True,
                          encoding="utf-8", errors="replace", cwd=str(REPO))


class ProfileSpecTests(unittest.TestCase):
    def test_every_profile_is_complete(self):
        keys = {"label", "note", "sizes", "lead_max_chars", "density_max",
                "table_rows_max", "matrix_rows_per_slide", "bullets_max_chars"}
        size_keys = set(dp.PROFILES[dp.DEFAULT_PROFILE]["sizes"])
        for name, spec in dp.PROFILES.items():
            with self.subTest(profile=name):
                self.assertEqual(keys - spec.keys(), set(), f"{name} 누락 필드")
                self.assertEqual(set(spec["sizes"]), size_keys, f"{name} 크기 항목 불일치")

    def test_presentation_is_larger_and_sparser_than_submission(self):
        pres = dp.get("presentation")
        detail = dp.get("detailed-submission")
        self.assertGreaterEqual(pres["sizes"]["body"], 18, "발표본 본문은 18pt 이상")
        for key in ("title", "lead", "body", "table"):
            self.assertGreater(pres["sizes"][key], detail["sizes"][key], key)
        self.assertLess(pres["density_max"], detail["density_max"])
        self.assertLess(pres["lead_max_chars"], detail["lead_max_chars"])
        self.assertLess(pres["matrix_rows_per_slide"], detail["matrix_rows_per_slide"])

    def test_executive_sits_between_the_two(self):
        ex, pres, detail = dp.get("executive-summary"), dp.get("presentation"), dp.get("detailed-submission")
        for key in ("body", "table", "title"):
            self.assertGreater(ex["sizes"][key], detail["sizes"][key], key)
            self.assertLess(ex["sizes"][key], pres["sizes"][key], key)

    def test_min_font_is_derived_not_hardcoded(self):
        """검사 하한이 생성기의 크기에서 나와야 프로파일을 바꿔도 어긋나지 않는다.

        본문과 소형 텍스트(표·도형 주석)의 하한은 서로 다르다 — 하나로 합치면 표 하한이
        본문에도 적용돼 본문을 표 크기까지 줄인 장표가 통과한다.
        """
        for name, spec in dp.PROFILES.items():
            with self.subTest(profile=name):
                self.assertEqual(dp.min_body_font(name), spec["sizes"]["body"] - 1)
                self.assertEqual(dp.min_table_font(name), spec["sizes"]["table"] - 1)
                self.assertLess(dp.min_table_font(name), dp.min_body_font(name), name)

    def test_unknown_profile_is_an_error(self):
        with self.assertRaises(ValueError):
            dp.get("executive")

    def test_stamp_round_trips(self):
        for name in dp.PROFILES:
            self.assertEqual(dp.from_stamp(dp.stamp(name)), name)
        for junk in ("", "presentation", "proposal-deck:nope", None, 3):
            self.assertIsNone(dp.from_stamp(junk), junk)


@unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
class BuilderProfileTests(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        self.addCleanup(self.tmp.cleanup)

    def _build(self, *extra: object, spec: Path | None = None) -> Path:
        out = self.dir / f"deck{len(list(self.dir.glob('*.pptx')))}.pptx"
        proc = run(BD, spec or FIX, "-o", out, *extra)
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        return out

    def _sizes(self, path: Path) -> set[float]:
        """모든 텍스트 크기(메타 도형 포함)."""
        prs = Presentation(str(path))
        return {r.font.size.pt for s in prs.slides for sh in s.shapes
                if sh.has_text_frame for p in sh.text_frame.paragraphs
                for r in p.runs if r.font.size is not None}

    def _body_min(self, path: Path) -> float:
        """검사기가 실제로 보는 본문 최소 폰트(캡션·헤더·푸터 제외)."""
        import deck_check
        prs = Presentation(str(path))
        found = [deck_check.min_font_pt(s)[0] for s in prs.slides]
        return min(v for v in found if v is not None)

    def test_profile_changes_actual_font_sizes(self):
        detail = self._sizes(self._build())
        pres = self._sizes(self._build("--profile", "presentation"))
        self.assertGreater(max(pres), max(detail))
        # 본문은 프로파일 하한 이상, 그리고 상세본보다 확실히 크다
        self.assertGreaterEqual(self._body_min(self._build("--profile", "presentation")),
                                dp.min_body_font("presentation"))
        self.assertGreater(self._body_min(self._build("--profile", "presentation")),
                           self._body_min(self._build()))

    def test_profile_is_stamped_into_the_file(self):
        for name in dp.PROFILES:
            with self.subTest(profile=name):
                out = self._build("--profile", name)
                self.assertEqual(Presentation(str(out)).core_properties.category, dp.stamp(name))

    def test_meta_output_profile_is_honoured(self):
        spec = json.loads(FIX.read_text(encoding="utf-8"))
        spec["meta"]["output_profile"] = "executive-summary"
        path = self.dir / "meta_profile.json"
        path.write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        out = self._build(spec=path)
        self.assertEqual(Presentation(str(out)).core_properties.category, dp.stamp("executive-summary"))

    def test_cli_profile_overrides_meta(self):
        spec = json.loads(FIX.read_text(encoding="utf-8"))
        spec["meta"]["output_profile"] = "executive-summary"
        path = self.dir / "meta_profile2.json"
        path.write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        out = self._build("--profile", "presentation", spec=path)
        self.assertEqual(Presentation(str(out)).core_properties.category, dp.stamp("presentation"))

    def test_unknown_profile_is_a_usage_error(self):
        proc = run(BD, FIX, "-o", self.dir / "x.pptx", "--profile", "huge")
        self.assertEqual(proc.returncode, 2)

    def test_presentation_splits_the_matrix_into_more_slides(self):
        detail = len(Presentation(str(self._build())).slides)
        pres = len(Presentation(str(self._build("--profile", "presentation"))).slides)
        self.assertGreater(pres, detail, "발표본은 조견표를 더 잘게 나눠야 한다")


@unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
class CheckerProfileTests(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        self.addCleanup(self.tmp.cleanup)

    def _build(self, profile: str | None = None) -> Path:
        out = self.dir / f"{profile or 'default'}.pptx"
        args = [BD, FIX, "-o", out] + (["--profile", profile] if profile else [])
        proc = run(*args)
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        return out

    def test_each_profile_passes_its_own_standard(self):
        for name in dp.PROFILES:
            with self.subTest(profile=name):
                proc = run(DC, self._build(name))
                self.assertEqual(proc.returncode, 0, proc.stdout)
                self.assertIn(dp.PROFILES[name]["label"], proc.stdout)

    def test_checker_reads_the_stamp_without_being_told(self):
        proc = run(DC, self._build("presentation"))
        self.assertIn("파일 표시", proc.stdout)
        self.assertIn(f"최소 폰트 {dp.min_body_font('presentation')}pt", proc.stdout)

    def test_submission_deck_fails_presentation_standard(self):
        proc = run(DC, self._build("detailed-submission"), "--profile", "presentation")
        self.assertEqual(proc.returncode, 1, proc.stdout)
        self.assertIn("최소 폰트", proc.stdout)
        self.assertIn("파일에 남은 표시", proc.stdout)

    def test_explicit_min_font_still_wins(self):
        proc = run(DC, self._build("presentation"), "--min-font", "40")
        self.assertEqual(proc.returncode, 1)
        self.assertIn("< 40", proc.stdout)

    def test_emit_render_records_the_profile(self):
        out = self.dir / "render.json"
        proc = run(DC, self._build("presentation"), "--emit-render", out)
        self.assertEqual(proc.returncode, 0, proc.stdout)
        self.assertEqual(json.loads(out.read_text(encoding="utf-8"))["output_profile"], "presentation")


class SharedConstantTests(unittest.TestCase):
    """생성기와 검사기가 같은 정의를 읽는지 — 드리프트 재발 방지."""

    def test_builder_and_checker_share_the_spec(self):
        sys.path.insert(0, str(DOC / "scripts"))
        import build_deck
        import deck_check
        self.assertIs(build_deck.deck_profiles, deck_check.deck_profiles)
        self.assertEqual(build_deck.LEAD_MAX_CHARS, deck_check.LEAD_MAX)
        self.assertEqual(build_deck.MATRIX_ROWS_PER_SLIDE,
                         dp.get(dp.DEFAULT_PROFILE)["matrix_rows_per_slide"])
        self.assertEqual(deck_check.TABLE_ROWS_MAX,
                         dp.get(dp.DEFAULT_PROFILE)["table_rows_max"])
        self.assertEqual(deck_check.DENSITY_MAX, dp.get(dp.DEFAULT_PROFILE)["density_max"])


if __name__ == "__main__":
    unittest.main()


@unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
class CheckerSubstanceTests(unittest.TestCase):
    """검사가 실제로 동작하는지 — '조용히 아무것도 안 하는 검사' 재발 방지."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        self.addCleanup(self.tmp.cleanup)
        sys.path.insert(0, str(DOC / "scripts"))

    def _build(self, spec: dict, *extra: object) -> Path:
        path = self.dir / "s.json"
        path.write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        out = self.dir / f"deck{len(list(self.dir.glob('*.pptx')))}.pptx"
        proc = run(BD, path, "-o", out, *extra)
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        return out

    def test_density_counts_table_cells(self):
        """표는 GraphicFrame이라 셀 텍스트가 밀도에서 통째로 빠져 있었다."""
        import deck_check
        rows = [{"id": f"R{i}", "text": "가상화 노드 증설과 이관 대상 산정" * 2, "support": "O 수용",
                 "response_loc": f"Ⅲ-{i} p.{i}", "note": "제조사 확약 첨부"} for i in range(1, 7)]
        deck = self._build({"meta": {"title": "t", "bidder": "b"},
                            "slides": [{"type": "matrix", "title": "조견표", "lead": "요구 대응 현황.",
                                        "rows": rows}]}, "--profile", "presentation")
        prs = Presentation(str(deck))
        table_slides = [s for s in prs.slides
                        if any(getattr(sh, "has_table", False) and sh.has_table
                               for sh in deck_check.iter_shapes(s.shapes))]
        self.assertTrue(table_slides)
        counted = max(
            sum(len(deck_check.shape_text(sh)) for sh in deck_check.iter_shapes(s.shapes)
                if sh.name not in deck_check.META_SHAPES)
            for s in table_slides)
        self.assertGreater(counted, 200, "표 셀이 밀도 계산에 들어가지 않는다")
        proc = run(DC, deck)
        self.assertIn("자 >", proc.stdout, "빽빽한 표 장표에 밀도 경고가 뜨지 않는다")

    def test_duplicate_shape_names_do_not_shift_the_density_pairing(self):
        """간트 마일스톤이 둘 이상이면 도형 이름이 겹친다 — 짝이 밀리면 안 된다."""
        import deck_check
        deck = self._build({"meta": {"title": "t", "bidder": "b"}, "slides": [
            {"type": "gantt", "title": "일정", "lead": "단계별 일정입니다.", "req_ids": ["R1"],
             "months": 4, "tasks": [{"name": "설계", "start": 1, "end": 2, "phase": True,
                                     "milestones": [{"at": 1, "label": "착수 보고"},
                                                    {"at": 2, "label": "설계 승인"}]}]}]})
        prs = Presentation(str(deck))
        slide = list(prs.slides)[0]
        pairs = [(sh.name, deck_check.shape_text(sh)) for sh in deck_check.iter_shapes(slide.shapes)]
        names = [n for n, _ in pairs]
        self.assertLess(len(set(names)), len(names), "이 테스트는 이름 중복 상황을 재현해야 한다")
        # 한 번 순회한 pairs가 정답이다. dict/list를 따로 만들어 zip하면 짝이 밀린다.
        zipped = sum(len(t) for n, t in zip(dict.fromkeys(names), [t for _, t in pairs]))
        correct = sum(len(t) for _, t in pairs)
        self.assertNotEqual(zipped, correct, "중복 이름에서 zip 방식은 실제로 어긋난다")
        proc = run(DC, deck)
        self.assertEqual(proc.returncode, 0, proc.stdout)

    def test_legend_is_checked_not_exempt(self):
        """구성도 범례는 작성자가 쓴 내용이다 — 폰트 하한 검사에서 빠지면 안 된다."""
        import deck_check
        self.assertNotIn("BODY_LEGEND", deck_check.META_SHAPES)
        deck = self._build({"meta": {"title": "t", "bidder": "b"}, "slides": [
            {"type": "zones", "title": "구성도", "lead": "계층을 분리했습니다.", "req_ids": ["R1"],
             "legend": ["실선: 데이터", "점선: 관리"],
             "zones": [{"title": "관리", "items": [{"title": "Control Plane"}]}]}]})
        prs = Presentation(str(deck))
        for slide in prs.slides:
            body_pt, table_pt = deck_check.min_font_pt(slide)
            # 범례는 소형 텍스트로 분류되므로 표 하한으로 잰다. 어느 쪽이든 검사 대상이며,
            # 하한 미만이면 차단된다(면제가 아니다).
            self.assertTrue(body_pt is not None or table_pt is not None)
            if body_pt is not None:
                self.assertGreaterEqual(body_pt, dp.min_body_font(dp.DEFAULT_PROFILE))
            if table_pt is not None:
                self.assertGreaterEqual(table_pt, dp.min_table_font(dp.DEFAULT_PROFILE))
        self.assertEqual(run(DC, deck).returncode, 0)

    def test_derived_floor_matches_what_the_builder_emits(self):
        """하한 공식을 되풀이하지 않고, 실제 생성물의 최소 본문 크기와 비교한다."""
        import deck_check
        for name in dp.PROFILES:
            with self.subTest(profile=name):
                out = self.dir / f"floor_{name}.pptx"
                proc = run(BD, FIX, "-o", out, "--profile", name)
                self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
                prs = Presentation(str(out))
                measured = [deck_check.min_font_pt(s) for s in prs.slides]
                body = [v[0] for v in measured if v[0] is not None]
                small = [v[1] for v in measured if v[1] is not None]
                self.assertGreaterEqual(min(body), dp.min_body_font(name),
                                        f"{name}: 생성물이 자기 본문 하한보다 작은 본문을 낸다")
                if small:
                    self.assertGreaterEqual(min(small), dp.min_table_font(name),
                                            f"{name}: 생성물이 자기 표 하한보다 작은 소형 텍스트를 낸다")


@unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
class StampFailClosedTests(unittest.TestCase):
    """규격 표시가 없거나 모르는 값이면 가장 느슨한 기준으로 조용히 통과시키지 않는다."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        self.addCleanup(self.tmp.cleanup)
        self.deck = self.dir / "p.pptx"
        proc = run(BD, FIX, "-o", self.deck, "--profile", "presentation")
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)

    def _restamp(self, value: str, name: str) -> Path:
        out = self.dir / name
        out.write_bytes(self.deck.read_bytes())
        prs = Presentation(str(out))
        prs.core_properties.category = value
        prs.save(str(out))
        return out

    def test_missing_stamp_blocks_at_submission_warns_at_draft(self):
        deck = self._restamp("", "nostamp.pptx")
        self.assertEqual(run(DC, deck, "--stage", "submission").returncode, 1)
        draft = run(DC, deck, "--stage", "draft")
        self.assertEqual(draft.returncode, 0, draft.stdout)
        self.assertIn("규격 표시가 없다", draft.stdout)

    def test_unknown_stamp_is_not_treated_as_missing(self):
        deck = self._restamp("proposal-deck:board-briefing", "future.pptx")
        proc = run(DC, deck, "--stage", "submission")
        self.assertEqual(proc.returncode, 1)
        self.assertIn("모르는 규격 표시", proc.stdout)

    def test_explicit_profile_clears_the_stamp_problem(self):
        deck = self._restamp("", "nostamp2.pptx")
        proc = run(DC, deck, "--profile", "presentation", "--stage", "submission")
        self.assertEqual(proc.returncode, 0, proc.stdout)
