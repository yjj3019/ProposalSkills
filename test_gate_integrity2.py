"""회귀 테스트 — 2026-09 재점검(R01~R08)에서 재현된 결함의 재발 방지.

이전 묶음(test_gate_integrity.py)이 --no-explain만 검사해 설명문 모순(R07)을 놓쳤다.
여기서는 설명 경로까지 함께 본다. 정상 대조군을 같이 둬서 과민 차단도 잡는다.
"""
from __future__ import annotations

import copy
import hashlib
import json
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

REPO = Path(__file__).resolve().parent
BEST = REPO / "skills/create-best-proposal"
DOC = REPO / "skills/create-proposal-document"
UG = BEST / "scripts/unified_gate.py"
QG = DOC / "scripts/quality_gate.py"
BD = DOC / "scripts/build_deck.py"
DC = DOC / "scripts/deck_check.py"
FIXTURES = BEST / "fixtures"

sys.path.insert(0, str(REPO))
sys.path.insert(0, str(REPO / "skills/create-winning-proposal/scripts"))
sys.path.insert(0, str(BEST / "scripts"))
import ooxml_fixtures as fixtures  # noqa: E402
import proposal_gate as pg  # noqa: E402
import build_audit_from_meta as bam  # noqa: E402

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:  # pragma: no cover
    HAS_PPTX = False


def run(*args: object, **kw) -> subprocess.CompletedProcess:
    return subprocess.run([sys.executable, *map(str, args)], capture_output=True, text=True,
                          encoding="utf-8", errors="replace", cwd=str(REPO), **kw)


def ready_audit() -> dict:
    return json.loads((FIXTURES / "audit_ready_financial.json").read_text(encoding="utf-8"))


class SubmissionBypassTests(unittest.TestCase):
    """R01 — 입력값 하나로 검증 의무를 취소할 수 없다."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        self.addCleanup(self.tmp.cleanup)
        self.doc = fixtures.pptx(self.dir / "final.pptx", raw={
            "ppt/slides/slide1.xml": "<a:p><a:r><a:t>정상 본문</a:t></a:r></a:p>"})
        self.digest = "sha256:" + hashlib.sha256(self.doc.read_bytes()).hexdigest()

    def _gate(self, data: dict, *extra: str) -> subprocess.CompletedProcess:
        p = self.dir / "audit.json"
        p.write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
        return run(UG, p, *extra)

    def _bound(self) -> dict:
        data = ready_audit()
        data["render"]["artifact_hash"] = data["package"]["artifact_hash"] = self.digest
        return data

    def test_bound_document_is_the_positive_control(self):
        proc = self._gate(self._bound(), "--doc", str(self.doc), "--no-explain")
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        self.assertIn("STATUS: SUBMISSION-READY", proc.stdout)

    def test_artifact_required_false_cannot_skip_verification(self):
        for value in (False, None):
            with self.subTest(artifact_required=value):
                data = self._bound()
                if value is None:
                    data.pop("artifact_required")
                else:
                    data["artifact_required"] = value
                data["render"].update(verified=False, artifact_hash="", tool="", evidence=[])
                data["package"]["artifact_hash"] = "review-record"
                proc = self._gate(data, "--doc", str(self.doc), "--no-explain")
                self.assertNotEqual(proc.returncode, 0, proc.stdout)
                self.assertNotIn("SUBMISSION-READY", proc.stdout)

    def test_visual_review_must_be_approved_by_a_person(self):
        data = self._bound()
        data["render"]["visual_review_approved"] = False
        self.assertTrue(any("visual review is not approved" in f for f in pg.evaluate(data)))
        data["render"].update(visual_review_approved=True, visual_reviewer="")
        self.assertTrue(any("named reviewer" in f for f in pg.evaluate(data)))

    def test_layout_check_flag_must_not_be_false(self):
        data = self._bound()
        data["render"]["layout_checked"] = False
        self.assertTrue(any("layout check is missing" in f for f in pg.evaluate(data)))


class PackageValidityTests(unittest.TestCase):
    """R02 — 열리지 않는 패키지를 제출 가능으로 판정하지 않는다."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        self.addCleanup(self.tmp.cleanup)

    def test_missing_content_types_is_rejected(self):
        p = self.dir / "half.pptx"
        with zipfile.ZipFile(p, "w") as z:
            z.writestr("ppt/presentation.xml", "<p:presentation/>")
            z.writestr("ppt/slides/slide1.xml", "<a:p/>")
        proc = run(QG, p)
        self.assertEqual(proc.returncode, 2, proc.stdout)
        self.assertIn("필수 파트 없음", proc.stdout + proc.stderr)

    def test_package_without_slides_is_rejected(self):
        p = fixtures.pptx(self.dir / "noslide.pptx")
        with zipfile.ZipFile(p, "a") as z:
            pass
        # 슬라이드 파트를 제거한 사본을 만든다
        stripped = self.dir / "stripped.pptx"
        with zipfile.ZipFile(p) as src, zipfile.ZipFile(stripped, "w") as dst:
            for item in src.infolist():
                if not item.filename.startswith("ppt/slides/"):
                    dst.writestr(item, src.read(item.filename))
        proc = run(QG, stripped)
        self.assertEqual(proc.returncode, 2, proc.stdout)
        self.assertIn("본문 파트가 없다", proc.stdout + proc.stderr)

    def test_broken_main_xml_is_rejected(self):
        p = self.dir / "broken.pptx"
        fixtures.pptx(p)
        rebuilt = self.dir / "rebuilt.pptx"
        with zipfile.ZipFile(p) as src, zipfile.ZipFile(rebuilt, "w") as dst:
            for item in src.infolist():
                data = src.read(item.filename)
                if item.filename == "ppt/presentation.xml":
                    data = b"<p:presentation><unclosed>"
                dst.writestr(item, data)
        proc = run(QG, rebuilt)
        self.assertEqual(proc.returncode, 2, proc.stdout)
        self.assertIn("XML 파싱 실패", proc.stdout + proc.stderr)

    @unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
    def test_real_pptx_passes(self):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1)).text_frame.text = "정상 본문"
        p = self.dir / "real.pptx"
        prs.save(str(p))
        self.assertEqual(run(QG, p).returncode, 0)

    @unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
    def test_fixture_packages_open_in_a_real_loader(self):
        """픽스처가 '게이트는 통과하는데 PowerPoint는 못 여는' 파일이 아닌지 확인한다."""
        p = fixtures.pptx(self.dir / "fixture.pptx", raw={
            "ppt/slides/slide1.xml": "<a:p><a:r><a:t>본문</a:t></a:r></a:p>"})
        from pptx.opc.serialized import PackageReader
        PackageReader(str(p))  # [Content_Types].xml·_rels 없으면 여기서 실패한다


class ExplainConsistencyTests(unittest.TestCase):
    """R07 — 라벨·본문·종료 코드가 한 판정에서 나온다(설명 경로 포함)."""

    def test_audit_only_explain_never_says_submittable(self):
        proc = run(UG, FIXTURES / "audit_ready_financial.json", "--audit-only")
        self.assertEqual(proc.returncode, 0, proc.stdout)
        self.assertIn("STATUS: AUDIT-VALID", proc.stdout)
        self.assertIn("## 게이트 결과: AUDIT-VALID", proc.stdout)
        for forbidden in ("게이트 결과: SUBMISSION-READY", "제출 가능."):
            self.assertNotIn(forbidden, proc.stdout)

    def test_draft_explain_never_says_submittable(self):
        data = ready_audit()
        data.update(mode="draft", artifact_required=False)
        data["checks"]["submission"] = False
        data["package"]["required"] = False
        data["submission"]["cleared"] = False
        with tempfile.TemporaryDirectory() as tmp:
            p = Path(tmp) / "draft.json"
            p.write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
            proc = run(UG, p, "--stage", "draft")
            self.assertEqual(proc.returncode, 0, proc.stdout)
            self.assertIn("STATUS: DRAFT-READY", proc.stdout)
            self.assertNotIn("제출 가능.", proc.stdout)


class SupportAliasTests(unittest.TestCase):
    """R06 — 미수용 별칭과 예외·응답 위치가 도구 사이에서 유지된다."""

    def test_every_unsupported_alias_is_treated_alike(self):
        for alias in ("X", "X 미수용", "미수용", "미지원", "✗", "부적합"):
            with self.subTest(alias=alias):
                data = ready_audit()
                data["requirements"][0].update(support=alias, fit="")
                self.assertTrue(any("approved but not met" in f for f in pg.evaluate(data)), alias)

    def test_accepted_codes_are_not_flagged(self):
        for alias in ("O", "O 수용", "부분", "조건부", "N/A", "해당없음", ""):
            with self.subTest(alias=alias):
                data = ready_audit()
                data["requirements"][0].update(support=alias, fit="STRONG")
                self.assertFalse(any("approved but not met" in f for f in pg.evaluate(data)), alias)

    def test_builder_keeps_exception_and_response_refs(self):
        meta = json.loads((FIXTURES / "meta_sample.json").read_text(encoding="utf-8"))
        meta["requirements"][0].update(
            support="X", fit="GAP", response_refs=["Ⅲ-1 p.12"],
            exception={"granted_by": "발주처 계약담당", "evidence": ["질의응답 회신 2026-08-14 3항"]})
        built = bam.build_audit(meta)["requirements"][0]
        self.assertEqual(built["response_refs"], ["Ⅲ-1 p.12"])
        self.assertEqual(built["exception"]["granted_by"], "발주처 계약담당")
        self.assertTrue(built["exception"]["evidence"])

    def test_builder_rejects_malformed_exception(self):
        meta = json.loads((FIXTURES / "meta_sample.json").read_text(encoding="utf-8"))
        meta["requirements"][0]["exception"] = "발주처가 허용함"
        with self.assertRaises(ValueError) as ctx:
            bam.build_audit(meta)
        self.assertIn("exception must be an object", str(ctx.exception))


class EnumTypeTests(unittest.TestCase):
    """R08 — 잘못된 타입은 크래시가 아니라 구조화된 오류다."""

    def test_no_type_error_for_any_enum_field(self):
        for value in ({}, [], 3, None, True):
            for path in (("mode",), ("bid_decision",), ("artifact_mode",)):
                with self.subTest(value=value, path=path):
                    data = ready_audit()
                    data[path[0]] = value
                    pg.validate_schema(data)
                    pg.evaluate(data)
            for holder, key in (("claims", "kind"), ("claims", "status"),
                                ("requirements", "state"), ("inputs", "class"),
                                ("defects", "severity")):
                with self.subTest(value=value, holder=holder, key=key):
                    data = ready_audit()
                    data.setdefault(holder, [{}])
                    if not data[holder]:
                        data[holder] = [{}]
                    data[holder][0][key] = value
                    pg.validate_schema(data)
                    pg.evaluate(data)
            with self.subTest(value=value, field="package.checks.metadata"):
                data = ready_audit()
                data["package"]["checks"]["metadata"] = value
                pg.validate_schema(data)
                pg.evaluate(data)


@unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
class DeckBuilderSafetyTests(unittest.TestCase):
    """R03·R05 — 템플릿 잔존 콘텐츠와 조용한 행 유실을 막는다."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        self.addCleanup(self.tmp.cleanup)

    def _spec(self, path: Path, spec: dict) -> Path:
        path.write_text(json.dumps(spec, ensure_ascii=False), encoding="utf-8")
        return path

    def test_template_with_existing_slides_is_rejected(self):
        tpl = Presentation()
        tpl.slide_width, tpl.slide_height = Inches(13.333), Inches(7.5)
        s = tpl.slides.add_slide(tpl.slide_layouts[6])
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = \
            "이전고객은행 제안서 37억원"
        tp = self.dir / "tpl.pptx"
        tpl.save(str(tp))
        spec = self._spec(self.dir / "s.json",
                          {"meta": {"title": "신규", "bidder": "제안사", "page_limit": 1},
                           "slides": [{"type": "cover"}]})
        out = self.dir / "o.pptx"
        proc = run(BD, spec, "-o", out, "--template", tp, "--strict")
        self.assertEqual(proc.returncode, 2, proc.stdout + proc.stderr)
        self.assertIn("이미 있다", proc.stdout + proc.stderr)
        self.assertFalse(out.exists(), "거부했는데 파일이 생성됐다")

    def test_empty_template_is_accepted(self):
        tpl = Presentation()
        tpl.slide_width, tpl.slide_height = Inches(13.333), Inches(7.5)
        tp = self.dir / "empty.pptx"
        tpl.save(str(tp))
        spec = self._spec(self.dir / "s2.json",
                          {"meta": {"title": "신규", "bidder": "제안사"},
                           "slides": [{"type": "cover"}]})
        out = self.dir / "o2.pptx"
        self.assertEqual(run(BD, spec, "-o", out, "--template", tp).returncode, 0)
        self.assertEqual(len(Presentation(str(out)).slides), 1)

    def test_invalid_rows_per_slide_is_a_usage_error(self):
        rows = [{"id": "R1", "text": "필수 요구", "support": "O", "response_loc": "p.1", "note": ""}]
        for bad in (-1, 0, "3", True, 1.5):
            with self.subTest(rows_per_slide=bad):
                spec = self._spec(self.dir / "m.json", {
                    "meta": {"title": "t", "bidder": "b"},
                    "slides": [{"type": "matrix", "title": "조견표", "lead": "요구 대응 현황.",
                                "rows_per_slide": bad, "rows": rows}]})
                out = self.dir / "m.pptx"
                out.unlink(missing_ok=True)
                proc = run(BD, spec, "-o", out, "--strict")
                self.assertEqual(proc.returncode, 2, f"{bad}: {proc.stdout}{proc.stderr}")

    def test_all_matrix_rows_survive_pagination(self):
        rows = [{"id": f"R{i}", "text": f"요구 {i}", "support": "O",
                 "response_loc": f"p.{i}", "note": ""} for i in range(1, 26)]
        for per in (1, 12, 13, 25, 40):
            with self.subTest(rows_per_slide=per):
                spec = self._spec(self.dir / "m2.json", {
                    "meta": {"title": "t", "bidder": "b"},
                    "slides": [{"type": "matrix", "title": "조견표", "lead": "요구 대응 현황.",
                                "rows_per_slide": per, "rows": rows}]})
                out = self.dir / "m2.pptx"
                out.unlink(missing_ok=True)
                proc = run(BD, spec, "-o", out, "--strict")
                self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
                seen = []
                for slide in Presentation(str(out)).slides:
                    for shape in slide.shapes:
                        if getattr(shape, "has_table", False) and shape.has_table:
                            seen += [r.cells[0].text for r in list(shape.table.rows)[1:]]
                self.assertEqual(seen, [r["id"] for r in rows], f"per={per}")

    def test_page_limit_counts_actual_slides(self):
        rows = [{"id": f"R{i}", "text": "요구", "support": "O", "response_loc": "p", "note": ""}
                for i in range(1, 26)]
        spec = self._spec(self.dir / "m3.json", {
            "meta": {"title": "t", "bidder": "b", "page_limit": 2},
            "slides": [{"type": "matrix", "title": "조견표", "lead": "요구 대응 현황.",
                        "rows_per_slide": 5, "rows": rows}]})
        out = self.dir / "m3.pptx"
        proc = run(BD, spec, "-o", out, "--strict")
        self.assertEqual(proc.returncode, 1)
        self.assertIn("페이지 제한", proc.stdout)
        self.assertIn("총 5장", proc.stdout)


@unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
class LayoutBoundsTests(unittest.TestCase):
    """R04 — 화면 밖·잘린 콘텐츠를 잡는다(정상 장표는 막지 않는다)."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        self.addCleanup(self.tmp.cleanup)

    def _deck(self, body_left_in: float) -> Path:
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for name, text, x, y in (("TITLE", "제목", 0.5, 0.6),
                                 ("LEAD", "결론을 먼저 씁니다.", 0.5, 1.2),
                                 ("BODY", "본문 R1", body_left_in, 2.0)):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(6), Inches(1))
            box.name = name
            box.text_frame.text = text
            box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        p = self.dir / f"deck_{body_left_in}.pptx"
        prs.save(str(p))
        return p

    def test_off_slide_body_is_blocked(self):
        proc = run(DC, self._deck(15.0), "--profile", "detailed-submission")
        self.assertEqual(proc.returncode, 1, proc.stdout)
        self.assertIn("화면 밖", proc.stdout)

    def test_inside_body_passes(self):
        # 외부에서 만든 덱은 규격 표시가 없으므로 어떤 기준으로 잴지 명시해야 한다.
        proc = run(DC, self._deck(0.5), "--profile", "detailed-submission")
        self.assertEqual(proc.returncode, 0, proc.stdout)

    def test_emit_render_separates_render_from_visual_review(self):
        out = self.dir / "render.json"
        proc = run(DC, self._deck(0.5), "--profile", "detailed-submission", "--emit-render", out)
        self.assertEqual(proc.returncode, 0, proc.stdout)
        block = json.loads(out.read_text(encoding="utf-8"))
        self.assertTrue(block["layout_checked"])
        self.assertFalse(block["visual_review_approved"], "육안 승인을 자동으로 참으로 두면 안 된다")
        self.assertIn("visual_reviewer", block)
        self.assertTrue(any("육안 검토 미완료" in e for e in block["evidence"]))


if __name__ == "__main__":
    unittest.main()
