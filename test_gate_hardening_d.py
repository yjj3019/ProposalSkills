"""8차 하드닝(D01–D09) 회귀 — 이번 배치에서 막은 허위 통과·허위 차단을 고정한다.

각 테스트는 "고치기 전에는 통과했다"를 재현한다. 검사 자체가 아니라 검사가 없던
자리를 지키는 것이 목적이므로, 규칙을 되풀이하지 않고 결과(차단/통과)로 확인한다.
"""
from __future__ import annotations

import json
import re
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

REPO = Path(__file__).resolve().parent
DOC = REPO / "skills/create-proposal-document"
BEST = REPO / "skills/create-best-proposal"
WIN = REPO / "skills/create-winning-proposal"
FIX = DOC / "fixtures/e2e-mini-rfp"
PY = sys.executable

sys.path.insert(0, str(WIN / "scripts"))
sys.path.insert(0, str(BEST / "scripts"))
sys.path.insert(0, str(DOC / "scripts"))
import proposal_gate as pg  # noqa: E402
import build_audit_from_meta as bam  # noqa: E402
import check_numbers as cn  # noqa: E402
import quality_gate as qg  # noqa: E402

sys.path.insert(0, str(REPO))
import ooxml_fixtures as fixtures  # noqa: E402
import install_skill  # noqa: E402

try:
    from pptx import Presentation  # noqa: F401
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def run(*args, **kw) -> subprocess.CompletedProcess:
    return subprocess.run([PY, *map(str, args)], capture_output=True, text=True,
                          encoding="utf-8", errors="replace", cwd=str(REPO), **kw)


class D01ConversionKeepsApprovalTests(unittest.TestCase):
    """meta→audit 변환이 검사·승인 기록을 버리면, 승인까지 마친 정상 경로가 차단된다."""

    def test_render_flags_survive_conversion(self):
        meta = json.loads((FIX / "meta.json").read_text(encoding="utf-8"))
        meta["render"] = {"verified": True, "render_succeeded": True, "layout_checked": True,
                          "visual_review_approved": True, "visual_reviewer": "QA",
                          "output_profile": "detailed-submission",
                          "artifact_hash": "sha256:" + "a" * 64, "tool": "deck_check",
                          "evidence": ["전 페이지 확인"]}
        audit = bam.build(meta, strict=False, warnings=[]) if hasattr(bam, "build") else None
        if audit is None:  # 공개 API 이름이 다르면 CLI로 확인한다
            with tempfile.TemporaryDirectory() as d:
                mp, ap = Path(d) / "m.json", Path(d) / "a.json"
                mp.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")
                self.assertEqual(run(BEST / "scripts/build_audit_from_meta.py", mp, "-o", ap).returncode, 0)
                audit = json.loads(ap.read_text(encoding="utf-8"))
        for flag in ("render_succeeded", "layout_checked", "visual_review_approved"):
            self.assertIs(audit["render"].get(flag), True, flag)
        self.assertEqual(audit["render"].get("visual_reviewer"), "QA")
        self.assertEqual(audit["render"].get("output_profile"), "detailed-submission")

    def test_conversion_does_not_invent_approval(self):
        meta = json.loads((FIX / "meta.json").read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as d:
            mp, ap = Path(d) / "m.json", Path(d) / "a.json"
            mp.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")
            run(BEST / "scripts/build_audit_from_meta.py", mp, "-o", ap)
            audit = json.loads(ap.read_text(encoding="utf-8"))
        self.assertIs(audit["render"].get("visual_review_approved"), False)


class D02RenderRecordTests(unittest.TestCase):
    """레이아웃 검사 기록은 '있으면 검사'가 아니라 필수다 — 지우면 요구가 사라졌다."""

    def _submission(self, **render):
        from test_proposal_gate import ready_data  # 공유 픽스처
        data = ready_data()
        data["render"].update(render)
        return pg.evaluate(data)

    def test_missing_layout_check_blocks(self):
        from test_proposal_gate import ready_data
        data = ready_data()
        data["render"].pop("layout_checked")
        self.assertTrue(any("layout check is missing" in f for f in pg.evaluate(data)))

    def test_verified_true_with_failed_render_is_a_contradiction(self):
        failures = self._submission(render_succeeded=False)
        self.assertTrue(any("contradicts render_succeeded=false" in f for f in failures))


class D04ArithmeticTests(unittest.TestCase):
    """원장 산술의 계산 불가·자기참조·금액 오차를 통과로 두지 않는다."""

    def test_non_finite_value_is_rejected(self):
        failures = pg.check_numbers([{"id": "N1", "label": "a", "value": float("inf"), "unit": "KRW"}])
        self.assertTrue(any("finite" in f for f in failures))

    def test_percent_needs_a_finite_amount_and_nonzero_base(self):
        base_zero = [{"id": "N1", "label": "모수", "value": 0, "unit": "KRW"},
                     {"id": "N2", "label": "비율", "value": 10, "unit": "%",
                      "percent_of": "N1", "amount": 5}]
        self.assertTrue(any("zero or non-numeric base" in f for f in pg.check_numbers(base_zero)))
        no_amount = [{"id": "N1", "label": "모수", "value": 100, "unit": "KRW"},
                     {"id": "N2", "label": "비율", "value": 10, "unit": "%", "percent_of": "N1"}]
        self.assertTrue(any("lacks a finite 'amount'" in f for f in pg.check_numbers(no_amount)))

    def test_currency_uses_absolute_tolerance(self):
        """1원 차이는 반올림이 아니라 오류다 — 상대 오차로는 큰 금액에서 조용히 통과했다."""
        entries = [{"id": "N1", "label": "합계", "value": 1_000_000_000, "unit": "KRW",
                    "components": ["N2", "N3"]},
                   {"id": "N2", "label": "a", "value": 600_000_000, "unit": "KRW"},
                   {"id": "N3", "label": "b", "value": 399_000_000, "unit": "KRW"}]
        self.assertTrue(any("합계가 맞지 않는다" in f for f in pg.check_numbers(entries)))

    def test_empty_ledger_in_submission_needs_a_reason(self):
        from test_proposal_gate import ready_data
        data = ready_data()
        data["numbers"] = []
        self.assertTrue(any("numbers ledger is empty" in f for f in pg.evaluate(data)))
        data["numbers_not_applicable"] = ["정성 평가 문서로 수치 없음(검토자 확인)"]
        self.assertFalse(any("numbers ledger is empty" in f for f in pg.evaluate(data)))


class D03NumberMatchingTests(unittest.TestCase):
    """같은 숫자라도 소수·단위·부호가 다르면 다른 값이다."""

    def test_decimal_and_prefix_do_not_count_as_a_match(self):
        self.assertFalse(cn._found("기간 37.5개월", "37", "개월"))
        self.assertFalse(cn._found("370원", "37", "원"))

    def test_unit_mismatch_does_not_count(self):
        self.assertFalse(cn._found("37개월 소요", "37", "원"))
        self.assertTrue(cn._found("총 37원", "37", "원"))

    def test_negative_value_does_not_satisfy_a_positive_entry(self):
        self.assertFalse(cn._found("-37원", "37", "원"))

    def test_korean_notation_still_matches(self):
        self.assertTrue(cn._found("사업비 37억원", "37억", "KRW"))
        self.assertTrue(cn._found("예산 3,700,000,000원", "3,700,000,000", "KRW"))

    def test_notes_only_value_is_reported_as_not_in_the_body(self):
        entries = [{"id": "N1", "label": "사업비", "value": 37, "unit": "원"}]
        items, _ = cn.compare(entries, "본문에는 수치가 없다", other="노트에 37원")
        self.assertTrue(any("비본문 영역에만" in i for i in items))


class D05PackageStructureTests(unittest.TestCase):
    """열리지 않는 패키지를 검사해 통과시키지 않는다."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        self.addCleanup(self.tmp.cleanup)

    def _broken(self, part: str) -> Path:
        import io
        import zipfile
        good = self.dir / "good.pptx"
        fixtures.pptx(good, {"ppt/slides/slide1.xml": "본문"})
        bad = self.dir / f"bad_{part.replace(chr(47), chr(95))}.pptx"
        buf = io.BytesIO()
        with zipfile.ZipFile(good) as src, zipfile.ZipFile(buf, "w") as out:
            for name in src.namelist():
                out.writestr(name, b"<p:sld><broken" if name.endswith(part) else src.read(name))
        bad.write_bytes(buf.getvalue())
        return bad

    def test_broken_slide_part_blocks(self):
        with self.assertRaises(KeyError):
            qg.extract_labeled_blocks(self._broken("slide1.xml"))

    def test_broken_relationships_block(self):
        with self.assertRaises(KeyError):
            qg.extract_labeled_blocks(self._broken("_rels/.rels"))

    def test_intact_package_still_reads(self):
        good = self.dir / "ok.pptx"
        fixtures.pptx(good, {"ppt/slides/slide1.xml": "본문 텍스트"})
        self.assertTrue(qg.extract_labeled_blocks(good))

    @unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
    def test_load_check_rejects_what_the_real_loader_rejects(self):
        synthetic = self.dir / "synthetic.pptx"
        fixtures.pptx(synthetic, {"ppt/slides/slide1.xml": "본문"})
        self.assertIsNotNone(qg.load_check(synthetic))


@unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
class D06LayoutInspectionTests(unittest.TestCase):
    """폰트 하한과 화면 밖 판정이 실제 장표를 정확히 재야 한다."""

    def test_body_and_small_text_have_separate_floors(self):
        import deck_profiles as dp
        for name in dp.PROFILES:
            self.assertLess(dp.min_table_font(name), dp.min_body_font(name), name)

    def test_paragraph_level_size_is_not_invisible_to_the_floor(self):
        import deck_check
        from pptx import Presentation as P
        from pptx.util import Inches, Pt
        prs = P()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.name = "BODY_TEXT"
        para = box.text_frame.paragraphs[0]
        para.text = "문단 수준에서 크기를 준 본문"
        para.font.size = Pt(6)          # run에는 크기가 없다
        body_pt, _ = deck_check.min_font_pt(slide)
        self.assertEqual(body_pt, 6.0)

    def test_group_children_are_measured_in_slide_coordinates(self):
        import deck_check
        out = Path(tempfile.mkdtemp()) / "g.pptx"
        proc = run(DOC / "scripts/build_deck.py", FIX / "slides.json", "-o", out)
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        from pptx import Presentation as P
        prs = P(str(out))
        for slide in prs.slides:
            self.assertEqual(deck_check.out_of_bounds(prs, slide), [])


class D07InstallerTests(unittest.TestCase):
    """설치기가 사용자 자료를 지우지 않는다."""

    def test_non_skill_directory_is_not_deleted(self):
        with tempfile.TemporaryDirectory() as d:
            root = Path(d)
            target = root / "create-best-proposal"
            target.mkdir()
            keep = target / "내 자료.md"
            keep.write_text("사용자 파일", encoding="utf-8")
            with self.assertRaises(SystemExit) as ctx:
                install_skill.install(root, "create-best-proposal")
            self.assertIn("Not empty", str(ctx.exception))
            self.assertTrue(keep.is_file())

    def test_force_replaces(self):
        with tempfile.TemporaryDirectory() as d:
            root = Path(d)
            (root / "create-best-proposal").mkdir()
            (root / "create-best-proposal" / "x.md").write_text("x", encoding="utf-8")
            target = install_skill.install(root, "create-best-proposal", force=True)
            self.assertTrue((target / "SKILL.md").is_file())


class D08ExplainTests(unittest.TestCase):
    """문서를 보지 않았으면 '해시 일치'라고 쓰지 않는다."""

    def test_audit_only_explain_does_not_claim_a_hash_match(self):
        text = pg.explain_markdown({"mode": "submission"}, [], [], label="SUBMISSION-READY",
                                   document_verified=False)
        self.assertNotIn("해시가 일치", text)
        self.assertIn("제출 승인이 아니다", text)

    def test_document_verified_explain_says_so(self):
        text = pg.explain_markdown({"mode": "submission"}, [], [], label="SUBMISSION-READY",
                                   document_verified=True)
        self.assertIn("해시가 일치", text)


class D09LedgerSubstanceTests(unittest.TestCase):
    """ID만 있는 껍데기 원장과, 검사에서 통째로 빠지던 분류를 막는다."""

    def _submission(self, **over):
        data = {"mode": "submission", "requirements": [], "claims": []}
        data.update(over)
        return pg.validate_schema(data)

    def test_requirement_without_text_is_rejected_in_submission(self):
        failures = self._submission(requirements=[{"id": "R1", "mandatory": True, "state": "approved"}])
        self.assertTrue(any("human-readable text" in f for f in failures))

    def test_draft_ledger_may_still_be_partial(self):
        data = {"mode": "draft", "requirements": [{"id": "R1", "state": "drafted"}], "claims": []}
        self.assertFalse(any("human-readable text" in f for f in pg.validate_schema(data)))

    def test_informational_claim_must_say_why_evidence_is_not_needed(self):
        failures = self._submission(claims=[{"id": "C1", "text": "제품 개요", "kind": "informational"}])
        self.assertTrue(any("lacks a rationale" in f for f in failures))
        ok = self._submission(claims=[{"id": "C1", "text": "제품 개요", "kind": "informational",
                                       "rationale": "사실 서술이며 성능 주장이 아니다"}])
        self.assertFalse(any("lacks a rationale" in f for f in ok))

    def test_malformed_unit_does_not_crash_the_gate(self):
        failures = pg.check_numbers([
            {"id": "N1", "label": "a", "value": 1, "unit": {"bad": True}},
            {"id": "N2", "label": "b", "value": 1, "unit": "KRW", "components": ["N1"]}])
        self.assertTrue(any("lacks a unit" in f for f in failures))


@unittest.skipUnless(HAS_PPTX, "python-pptx 없음")
class NormalPathTests(unittest.TestCase):
    """정상 경로가 실제로 끝까지 간다 — 강화가 정상 산출물을 막지 않는지 확인한다.

    승인된 입력·원장 → PPTX 생성 → 배치·수치·패키지 검사 → 그 파일에 대한 승인 기록 →
    meta→audit 변환 → 같은 파일에 대한 최종 판정.
    """

    def test_approved_inputs_reach_submission_ready_on_the_same_file(self):
        import hashlib
        tmp = tempfile.TemporaryDirectory()
        self.addCleanup(tmp.cleanup)
        out = Path(tmp.name)
        pptx = out / "final.pptx"
        # 골든 픽스처는 일부러 미완성(○○·NEEDS INPUT)이다. 정상 경로 검증에는
        # 그 자리를 채운 사본을 쓴다 — 채우지 않으면 품질 게이트에서 정당하게 막힌다.
        spec = json.loads((FIX / "slides.json").read_text(encoding="utf-8"))

        def resolve(node):
            if isinstance(node, str):
                text = re.sub(r"\[NEEDS INPUT:[^\]]*\]", "2026-12-31", node)
                return text.replace("○○", "가상")
            if isinstance(node, list):
                return [resolve(v) for v in node]
            if isinstance(node, dict):
                return {k: resolve(v) for k, v in node.items()}
            return node

        resolved = out / "slides_resolved.json"
        resolved.write_text(json.dumps(resolve(spec), ensure_ascii=False), encoding="utf-8")
        self.assertEqual(run(DOC / "scripts/build_deck.py", resolved,
                             "-o", pptx, "--strict").returncode, 0)
        # 1) 배치 검사 — 실제 파일을 연다
        check = run(DOC / "scripts/deck_check.py", pptx, "--stage", "draft")
        self.assertEqual(check.returncode, 0, check.stdout)
        # 2) 수치 대조 — 원장이 본문에 실제로 있는지
        meta = json.loads((FIX / "meta.json").read_text(encoding="utf-8"))
        ledger = out / "numbers.json"
        ledger.write_text(json.dumps(meta["numbers"], ensure_ascii=False), encoding="utf-8")
        nums = run(DOC / "scripts/check_numbers.py", pptx, "--numbers", ledger)
        self.assertEqual(nums.returncode, 0, nums.stdout)
        # 3) 승인 기록은 이 파일의 해시에 묶는다
        digest = "sha256:" + hashlib.sha256(pptx.read_bytes()).hexdigest()
        meta.update(mode="submission", bid_decision="bid", bid_conditions=[], unresolved_tokens=[])
        meta["eligibility"][0]["met"] = True
        for r in meta["requirements"]:
            r.update(state="approved", evidence_refs=["제안서 본문"], reviewer="검토자")
            r.setdefault("text", f"{r['id']} 요구 내용")
        for c in meta.get("claims", []):
            c.update(status="supported", owner_approved=True, evidence_refs=["시험 보고서"])
            c.setdefault("text", f"{c['id']} 주장 내용")
        for a in meta["attachments"]:
            a["present"] = True
        for i in meta["inputs"]:
            i["status"] = "closed"
        meta["checks"] = {"consistency": True, "arithmetic": True, "submission": True}
        meta["artifact_required"] = True
        meta["render"] = {"verified": True, "render_succeeded": True, "layout_checked": True,
                          "visual_review_approved": True, "visual_reviewer": "QA",
                          "output_profile": "detailed-submission", "artifact_hash": digest,
                          "tool": "deck_check", "evidence": ["전 페이지 육안 확인"]}
        meta["package"] = {"required": True, "inspected": True, "artifact_hash": digest,
                           "tool": "ooxml-check", "reviewer": "QA",
                           "checks": {k: "pass" for k in
                                      ("metadata", "notes", "comments", "hidden-content",
                                       "external-links", "stale-customer-data", "price-leakage")}
                           | {"embedded-files": "not-applicable", "macros": "not-applicable"}}
        meta["submission"] = {"cleared": True, "deadline": "2099-12-31T17:00:00+09:00",
                              "rehearsal_evidence": ["시험 업로드 완료"],
                              "receipt_plan": "접수증 보관", "receipt_evidence": []}
        mp, ap = out / "meta_final.json", out / "audit_final.json"
        mp.write_text(json.dumps(meta, ensure_ascii=False), encoding="utf-8")
        conv = run(BEST / "scripts/build_audit_from_meta.py", mp, "-o", ap)
        self.assertEqual(conv.returncode, 0, conv.stdout + conv.stderr)
        gate = run(BEST / "scripts/unified_gate.py", ap, "--doc", pptx, "--no-explain")
        self.assertEqual(gate.returncode, 0, gate.stdout + gate.stderr)
        self.assertIn("STATUS: SUBMISSION-READY", gate.stdout)


if __name__ == "__main__":
    unittest.main()
